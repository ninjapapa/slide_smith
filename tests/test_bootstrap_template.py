import json
import tempfile
from pathlib import Path

from pptx import Presentation

from slide_smith.template_bootstrapper import bootstrap_template
from slide_smith.template_validator import validate_template


def _make_temp_pptx() -> str:
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    tmpdir = tempfile.mkdtemp(prefix="slide-smith-")
    p = Path(tmpdir) / "seed.pptx"
    prs.save(str(p))
    return str(p)


def test_bootstrap_template_writes_package_and_validates() -> None:
    seed = _make_temp_pptx()
    out_dir = tempfile.mkdtemp(prefix="slide-smith-out-")

    res = bootstrap_template(pptx_path=seed, template_id="t1", out_dir=out_dir)

    tdir = Path(res.template_dir)
    assert (tdir / "template.pptx").exists()
    assert (tdir / "template.json").exists()

    spec = json.loads((tdir / "template.json").read_text())
    assert spec["template_id"] == "t1"
    assert isinstance(spec["archetypes"], list)
    assert len(spec["archetypes"]) > 0

    v = validate_template("t1", templates_dir=out_dir)
    assert v.ok


def test_bootstrap_template_overwrite_behavior() -> None:
    seed = _make_temp_pptx()
    out_dir = tempfile.mkdtemp(prefix="slide-smith-out-")

    _ = bootstrap_template(pptx_path=seed, template_id="t2", out_dir=out_dir)

    # Without overwrite should fail.
    try:
        _ = bootstrap_template(pptx_path=seed, template_id="t2", out_dir=out_dir)
        assert False, "expected bootstrap to fail without overwrite"
    except Exception:
        pass

    # With overwrite should succeed.
    _ = bootstrap_template(pptx_path=seed, template_id="t2", out_dir=out_dir, overwrite=True)
