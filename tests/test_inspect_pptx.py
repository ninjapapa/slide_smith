import json
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation


def _make_temp_pptx() -> str:
    prs = Presentation()
    # Add a slide to ensure the file is a valid pptx and has default layouts.
    prs.slides.add_slide(prs.slide_layouts[0])

    tmpdir = tempfile.mkdtemp(prefix="slide-smith-")
    path = Path(tmpdir) / "example.pptx"
    prs.save(str(path))
    return str(path)


def test_inspect_pptx_outputs_layouts_and_placeholders() -> None:
    repo = Path(__file__).resolve().parents[1]
    pptx_path = _make_temp_pptx()

    cmd = [
        str(repo / ".venv" / "bin" / "python"),
        "-m",
        "slide_smith.cli",
        "inspect-pptx",
        "--pptx",
        pptx_path,
    ]
    out = subprocess.check_output(cmd, cwd=str(repo))
    data = json.loads(out.decode("utf-8"))

    assert data["pptx"].endswith("example.pptx")
    assert "slide_size" in data
    assert "layouts" in data
    assert isinstance(data["layouts"], list)
    assert len(data["layouts"]) > 0

    first = data["layouts"][0]
    assert "name" in first
    assert "index" in first
    assert "placeholders" in first
    assert isinstance(first["placeholders"], list)

    for layout in data["layouts"]:
        for ph in layout.get("placeholders", []):
            assert isinstance(ph["idx"], int)
            assert isinstance(ph["ph_type"], str)


def test_inspect_pptx_is_deterministic() -> None:
    repo = Path(__file__).resolve().parents[1]
    pptx_path = _make_temp_pptx()

    cmd = [
        str(repo / ".venv" / "bin" / "python"),
        "-m",
        "slide_smith.cli",
        "inspect-pptx",
        "--pptx",
        pptx_path,
    ]

    out1 = subprocess.check_output(cmd, cwd=str(repo)).decode("utf-8")
    out2 = subprocess.check_output(cmd, cwd=str(repo)).decode("utf-8")
    assert out1 == out2
