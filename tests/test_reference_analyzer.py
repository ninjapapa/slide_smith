from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation

from slide_smith.reference_analyzer import analyze_reference


def test_analyze_reference_emits_layout_placeholders_with_bbox(tmp_path: Path) -> None:
    # Create a minimal PPTX using python-pptx default template.
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])

    pptx_path = tmp_path / "ref.pptx"
    prs.save(pptx_path)

    res = analyze_reference(str(pptx_path))
    profile = res.style_profile

    assert profile["version"] == "1"
    assert profile["reference"]["sha256"]
    assert profile["reference"]["slideSize"]["widthEmu"] > 0
    assert profile["reference"]["slideSize"]["heightEmu"] > 0

    layouts = profile["layouts"]
    assert isinstance(layouts, list)
    assert len(layouts) > 0

    # Ensure at least one placeholder has bbox keys.
    any_bbox = False
    for layout in layouts:
        assert "layoutId" in layout and isinstance(layout["layoutId"], str) and layout["layoutId"].startswith("layout:")
        for ph in layout.get("placeholders", []):
            bbox = ph.get("bbox")
            if isinstance(bbox, dict) and set(["x", "y", "w", "h"]).issubset(bbox.keys()):
                any_bbox = True
                break
        if any_bbox:
            break

    assert any_bbox, "Expected at least one placeholder bbox in StyleProfile"


def test_analyze_reference_is_json_serializable(tmp_path: Path) -> None:
    prs = Presentation()
    pptx_path = tmp_path / "ref.pptx"
    prs.save(pptx_path)

    res = analyze_reference(str(pptx_path))
    json.dumps(res.style_profile, sort_keys=True)
