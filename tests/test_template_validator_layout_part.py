from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation

from slide_smith.openxml_layouts import inspect_openxml_layouts
from slide_smith.template_validator import validate_template


def test_validate_template_allows_layout_part_match_when_layout_name_wrong(tmp_path: Path) -> None:
    templates_root = tmp_path / "templates"
    tdir = templates_root / "t2"
    tdir.mkdir(parents=True)

    prs = Presentation()
    prs.save(tdir / "template.pptx")

    raw = inspect_openxml_layouts(str(tdir / "template.pptx"))
    assert raw.layouts
    part0 = raw.layouts[0]["part"]
    phs = raw.layouts[0].get("placeholders") or []
    if not phs:
        # If this particular layout has no placeholders, pick another.
        for l in raw.layouts:
            if l.get("placeholders"):
                part0 = l["part"]
                phs = l.get("placeholders") or []
                break
    assert phs, "expected at least one placeholder in some slide layout"
    idx0 = int(phs[0]["idx"])

    spec = {
        "version": "1.0",
        "meta": {"template": "t2"},
        "archetypes": [
            {
                "id": "a1",
                "layout": "THIS NAME DOES NOT EXIST",
                "layout_part": part0,
                "slots": [
                    {"name": "title", "type": "text", "placeholder_idx": idx0, "required": False},
                ],
            }
        ],
    }
    (tdir / "template.json").write_text(json.dumps(spec), encoding="utf-8")

    res = validate_template("t2", templates_dir=str(templates_root), profile="structural")
    assert res.ok, f"expected ok; got: {res.errors}"
