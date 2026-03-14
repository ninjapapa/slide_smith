from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation

from slide_smith.reference_analyzer import analyze_reference
from slide_smith.exemplar_compiler import compile_exemplar
from slide_smith.exemplar_renderer import render_exemplar


def test_render_exemplar_renders_text_slides(tmp_path: Path) -> None:
    # Build a minimal reference deck using python-pptx defaults.
    ref = Presentation()
    ref.slides.add_slide(ref.slide_layouts[0])
    ref_path = tmp_path / "ref.pptx"
    ref.save(ref_path)

    profile = analyze_reference(str(ref_path)).style_profile

    # Build a slide_plan that avoids images.
    slide_plan = {
        "version": "1",
        "slides": [
            {"intent": "title", "fields": {"title": "Hello"}},
            {"intent": "bullets", "fields": {"title": "Agenda", "bullets": ["A", "B"]}},
        ],
    }

    spec = compile_exemplar(slide_plan=slide_plan, style_profile=profile).deck_spec

    out_path = tmp_path / "out.pptx"
    res = render_exemplar(
        reference_pptx=str(ref_path),
        style_profile=profile,
        deck_spec=spec,
        output_pptx=str(out_path),
    )

    assert Path(res.output_pptx).exists()

    # Verify output slide count and some text presence.
    prs = Presentation(res.output_pptx)
    assert len(prs.slides) == 2

    texts = []
    for s in prs.slides:
        for sh in s.shapes:
            if getattr(sh, "has_text_frame", False):
                t = (sh.text_frame.text or "").strip()
                if t:
                    texts.append(t)

    joined = "\n".join(texts)
    assert "Hello" in joined
    assert "Agenda" in joined


def test_render_exemplar_rejects_unknown_layout(tmp_path: Path) -> None:
    ref = Presentation()
    ref_path = tmp_path / "ref.pptx"
    ref.save(ref_path)

    profile = analyze_reference(str(ref_path)).style_profile

    deck_spec = {
        "version": "1",
        "reference": {"path": str(ref_path), "sha256": "x"},
        "slides": [{"layoutId": "layout:does-not-exist", "fills": []}],
    }

    try:
        render_exemplar(
            reference_pptx=str(ref_path),
            style_profile=profile,
            deck_spec=deck_spec,
            output_pptx=str(tmp_path / "out.pptx"),
        )
        assert False, "expected error"
    except ValueError as e:
        assert "Unknown layoutId" in str(e)
