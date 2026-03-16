from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation

from slide_smith.editor import delete_slide_from_deck, list_slides_in_deck, update_slide_in_deck


def _make_deck(tmp_path: Path) -> Path:
    from slide_smith.renderer import render_deck
    from slide_smith.template_loader import load_template_spec

    # minimal valid png
    img = tmp_path / "demo.png"
    img.write_bytes(
        bytes.fromhex(
            "89504e470d0a1a0a0000000d4948445200000001000000010804000000b51c0c020000000b4944415478da63fcff1f0002eb01f56c124babaa0000000049454e44ae426082"
        )
    )

    deck = {
        "slides": [
            {"layout_id": "title", "title": "T", "subtitle": "S"},
            {
                "layout_id": "text_with_image",
                "title": "I",
                "body": "B",
                "image": {"path": "demo.png", "alt": "a"},
            },
        ]
    }
    out = tmp_path / "out.pptx"
    render_deck(deck, load_template_spec("default"), "default", str(out), base_dir=str(tmp_path))
    return out


def test_list_and_delete_slide(tmp_path: Path) -> None:
    deck = _make_deck(tmp_path)
    items = list_slides_in_deck(str(deck))
    assert [x["index"] for x in items] == [0, 1]

    delete_slide_from_deck(str(deck), 0)
    prs = Presentation(str(deck))
    assert len(prs.slides) == 1


def test_update_slide_notes_and_image(tmp_path: Path) -> None:
    deck = _make_deck(tmp_path)

    # new image
    new_img = tmp_path / "new.png"
    new_img.write_bytes(
        bytes.fromhex(
            "89504e470d0a1a0a0000000d4948445200000001000000010804000000b51c0c020000000b4944415478da63fcff1f0002eb01f56c124babaa0000000049454e44ae426082"
        )
    )

    patch = tmp_path / "patch.json"
    patch.write_text(json.dumps({"notes": "hello", "image": {"path": "new.png", "alt": "x"}}))

    update_slide_in_deck(str(deck), 1, str(patch))
    prs = Presentation(str(deck))
    assert prs.slides[1].notes_slide.notes_text_frame.text == "hello"
