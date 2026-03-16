from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation

from slide_smith.commands.edit_ops import handle_update_slide
from slide_smith.renderer import render_deck
from slide_smith.template_loader import load_template_spec


def _base_deck(tmp_path: Path) -> Path:
    out = tmp_path / "base.pptx"
    render_deck(
        {"slides": [{"layout_id": "title", "title": "Demo", "subtitle": "Sub"}]},
        load_template_spec("default"),
        "default",
        str(out),
        base_dir=str(tmp_path),
    )
    return out


def test_update_slide_body_uses_known_placeholder_when_available(tmp_path: Path) -> None:
    deck = _base_deck(tmp_path)
    patch = tmp_path / "patch.json"
    patch.write_text(json.dumps({"body": "Body written into known placeholder"}))

    code, out = handle_update_slide(deck=str(deck), index=0, input_path=str(patch))
    assert code == 0
    assert "slide updated" in out

    prs = Presentation(str(deck))
    assert prs.slides[0].placeholders[1].text == "Body written into known placeholder"


def test_update_slide_invalid_index_returns_clear_error(tmp_path: Path) -> None:
    deck = _base_deck(tmp_path)
    patch = tmp_path / "patch.json"
    patch.write_text(json.dumps({"title": "Nope"}))

    code, out = handle_update_slide(deck=str(deck), index=5, input_path=str(patch))
    assert code == 1
    assert "Update-slide failed" in out
    assert "out of range" in out
