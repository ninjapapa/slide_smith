from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from slide_smith.reference_analyzer import analyze_reference
from slide_smith.exemplar_compiler import compile_exemplar
from slide_smith.exemplar_renderer import render_exemplar
from slide_smith.exemplar_validator import validate_exemplar


def test_validate_exemplar_happy_path(tmp_path: Path) -> None:
    # Build reference
    ref = Presentation()
    ref.slides.add_slide(ref.slide_layouts[0])
    ref_path = tmp_path / "ref.pptx"
    ref.save(ref_path)

    profile = analyze_reference(str(ref_path)).style_profile

    # Build spec and render
    slide_plan = {
        "version": "1",
        "slides": [
            {"intent": "title", "fields": {"title": "Hello"}},
            {"intent": "bullets", "fields": {"title": "Agenda", "bullets": ["A", "B"]}},
        ],
    }
    spec = compile_exemplar(slide_plan=slide_plan, style_profile=profile).deck_spec

    out_path = tmp_path / "out.pptx"
    render_exemplar(
        reference_pptx=str(ref_path),
        style_profile=profile,
        deck_spec=spec,
        output_pptx=str(out_path),
    )

    res = validate_exemplar(reference_pptx=str(ref_path), pptx=str(out_path), style_profile=profile)
    assert res.ok, f"expected ok, got errors: {res.errors}"
    assert res.errors == []


def test_validate_exemplar_flags_extra_shape(tmp_path: Path) -> None:
    ref = Presentation()
    ref.slides.add_slide(ref.slide_layouts[0])
    ref_path = tmp_path / "ref.pptx"
    ref.save(ref_path)

    profile = analyze_reference(str(ref_path)).style_profile

    # Make an output deck that adds an extra textbox shape (not a placeholder)
    prs = Presentation(str(ref_path))
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.add_textbox(0, 0, 1000000, 1000000).text_frame.text = "EXTRA"
    out_path = tmp_path / "out.pptx"
    prs.save(out_path)

    res = validate_exemplar(reference_pptx=str(ref_path), pptx=str(out_path), style_profile=profile)
    assert not res.ok
    assert any("unexpected shape" in e for e in res.errors)
