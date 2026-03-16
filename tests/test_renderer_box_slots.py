from __future__ import annotations

import base64
import json
from pathlib import Path

from pptx import Presentation

from slide_smith.renderer import render_deck


PNG_1X1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO7Wk0kAAAAASUVORK5CYII="
)


def test_render_box_based_text_with_image(tmp_path: Path):
    # Build a minimal template package on the fly.
    template_id = "t"
    templates_dir = tmp_path / "templates"
    tdir = templates_dir / template_id
    tdir.mkdir(parents=True)

    prs = Presentation()
    layout = prs.slide_layouts[6]
    layout_name = layout.name

    pptx_path = tdir / "template.pptx"
    prs.save(str(pptx_path))

    template_spec = {
        "template_id": template_id,
        "name": "tmp",
        "version": "0.0",
        "deck": {"aspect_ratio": "16:9", "supported_layout_ids": ["text_with_image"]},
        "archetypes": [
            {
                "id": "text_with_image",
                "layout": layout_name,
                "slots": [
                    {"name": "title", "type": "text", "required": True, "box": {"units": "relative", "x": 0.05, "y": 0.05, "w": 0.9, "h": 0.12}},
                    {"name": "image", "type": "image", "required": True, "box": {"units": "relative", "x": 0.55, "y": 0.25, "w": 0.4, "h": 0.6}},
                    {"name": "body", "type": "text", "required": True, "box": {"units": "relative", "x": 0.05, "y": 0.25, "w": 0.45, "h": 0.6}},
                ],
            }
        ],
        "styles": {},
    }

    (tdir / "template.json").write_text(json.dumps(template_spec))

    img_path = tmp_path / "img.png"
    img_path.write_bytes(PNG_1X1)

    deck_spec = {
        "slides": [
            {
                "layout_id": "text_with_image",
                "title": "Hello",
                "body": "World",
                "image": {"path": str(img_path)},
            }
        ]
    }

    out = tmp_path / "out.pptx"
    render_deck(deck_spec, template_spec, template_id, str(out), base_dir=str(tmp_path), templates_dir=str(templates_dir))

    rendered = Presentation(str(out))
    assert len(rendered.slides) == 1
    slide = rendered.slides[0]

    # We should have at least: 2 textboxes + 1 picture, regardless of the base layout.
    text_like = [s for s in slide.shapes if getattr(s, "has_text_frame", False)]
    assert any("Hello" in (s.text_frame.text or "") for s in text_like)
    assert any("World" in (s.text_frame.text or "") for s in text_like)

    # Picture shapes in python-pptx have .image
    pics = [s for s in slide.shapes if getattr(s, "image", None) is not None]
    assert len(pics) >= 1
