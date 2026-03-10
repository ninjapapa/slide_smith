import json
import tempfile
from pathlib import Path

from pptx import Presentation

from slide_smith.renderer import render_deck
from slide_smith.template_bootstrapper import bootstrap_template
from slide_smith.template_loader import load_template_spec
from slide_smith.template_validator import validate_template


def _make_seed_pptx() -> str:
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    tmpdir = tempfile.mkdtemp(prefix="slide-smith-")
    p = Path(tmpdir) / "seed.pptx"
    prs.save(str(p))
    return str(p)


def test_end_to_end_bootstrap_validate_render() -> None:
    seed = _make_seed_pptx()
    out_dir = tempfile.mkdtemp(prefix="slide-smith-out-")

    # 1) bootstrap
    _ = bootstrap_template(pptx_path=seed, template_id="t_e2e", out_dir=out_dir)

    # 2) validate
    v = validate_template("t_e2e", templates_dir=out_dir)
    assert v.ok

    # 3) render a dummy deck using an implemented archetype.
    # Bootstrapped archetypes are layout__... and are not rendered by the current renderer.
    # For v1, the end-to-end loop we care about is: bootstrap + validate + render using a known archetype.
    template_spec = load_template_spec("t_e2e", templates_dir=out_dir)

    # Add a minimal mapping so we can render using the existing "title" renderer.
    first_layout = template_spec["archetypes"][0]["layout"]
    template_spec["archetypes"].append(
        {
            "id": "title",
            "description": "Title slide (mapped for test)",
            "layout": first_layout,
            "slots": [
                {"name": "title", "type": "text", "required": False, "placeholder_idx": 0},
                {"name": "subtitle", "type": "text", "required": False, "placeholder_idx": 1},
            ],
        }
    )

    deck_spec = {
        "version": "1.0",
        "meta": {"title": "e2e"},
        "slides": [
            {"archetype": "title", "title": "Hello", "subtitle": "World"},
        ],
    }

    out_path = Path(tempfile.mkdtemp(prefix="slide-smith-out-pptx-") ) / "out.pptx"
    rendered = render_deck(
        deck_spec,
        template_spec,
        "t_e2e",
        str(out_path),
        base_dir=str(Path(seed).parent),
        templates_dir=out_dir,
    )
    assert Path(rendered).exists()
    assert Path(rendered).stat().st_size > 0
