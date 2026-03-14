from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from slide_smith.openxml_layouts import inspect_openxml_layouts
from slide_smith.reference_analyzer import analyze_reference


def test_openxml_layouts_enumerates_slide_layout_parts(tmp_path: Path) -> None:
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    p = tmp_path / "ref.pptx"
    prs.save(p)

    raw = inspect_openxml_layouts(str(p))
    assert raw.layouts, "expected at least one slideLayout part"

    # Basic contract: each layout has part path
    assert all(l.get("part", "").startswith("ppt/slideLayouts/") for l in raw.layouts)


def test_analyze_reference_raw_mode_runs(tmp_path: Path) -> None:
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    p = tmp_path / "ref.pptx"
    prs.save(p)

    prof = analyze_reference(str(p), mode="raw").style_profile
    assert prof["layouts"], "expected layouts in raw mode"
    assert prof["reference"]["slideSize"]["widthEmu"] > 0
