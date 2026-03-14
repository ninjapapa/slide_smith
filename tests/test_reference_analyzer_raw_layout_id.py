from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from slide_smith.reference_analyzer import analyze_reference


def test_analyze_reference_raw_layout_ids_use_part_name(tmp_path: Path) -> None:
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    p = tmp_path / "ref.pptx"
    prs.save(p)

    prof = analyze_reference(str(p), mode="raw").style_profile
    lids = [l.get("layoutId", "") for l in prof.get("layouts", [])]
    assert lids
    # Should use slideLayout stem rather than numeric index.
    assert any("layout:slidelayout" in lid.lower() for lid in lids)
