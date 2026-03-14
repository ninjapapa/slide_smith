from __future__ import annotations

import zipfile
from pathlib import Path

from pptx import Presentation

from slide_smith.potx_converter import convert_potx_to_pptx


def _make_fake_potx_from_pptx(pptx_in: Path, potx_out: Path) -> None:
    """Create a POTX-like package by rewriting [Content_Types].xml override."""

    with zipfile.ZipFile(pptx_in, "r") as zin:
        with zipfile.ZipFile(potx_out, "w", compression=zipfile.ZIP_DEFLATED) as zout:
            for info in zin.infolist():
                data = zin.read(info.filename)
                if info.filename == "[Content_Types].xml":
                    data = data.replace(
                        b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
                        b"application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
                    )
                zout.writestr(info, data)


def test_convert_potx_to_pptx_makes_python_pptx_loadable(tmp_path: Path) -> None:
    # Create a minimal pptx.
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    pptx = tmp_path / "in.pptx"
    prs.save(pptx)

    # Rewrite to potx-like; python-pptx should reject.
    potx = tmp_path / "in.potx"
    _make_fake_potx_from_pptx(pptx, potx)

    rejected = False
    try:
        Presentation(str(potx))
    except Exception:
        rejected = True
    assert rejected, "expected python-pptx to reject potx content-type"

    out = tmp_path / "out.pptx"
    convert_potx_to_pptx(potx_path=str(potx), pptx_path=str(out), overwrite=True)

    # Should load now.
    prs2 = Presentation(str(out))
    assert prs2.slide_width > 0
