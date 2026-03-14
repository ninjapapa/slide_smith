from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation

from slide_smith.openxml_layouts import inspect_openxml_layouts
from slide_smith.renderer import render_deck


def test_renderer_can_resolve_layout_by_part_when_name_is_wrong(tmp_path: Path) -> None:
    # Build a template presentation.
    prs = Presentation()
    pptx_path = tmp_path / "template.pptx"
    prs.save(pptx_path)

    raw = inspect_openxml_layouts(str(pptx_path))
    assert raw.layouts

    # Pick a layout that has placeholders with a non-negative idx.
    chosen = None
    idx0 = None
    for l in raw.layouts:
        phs = l.get("placeholders") or []
        for ph in phs:
            try:
                idx = int(ph.get("idx"))
            except Exception:
                continue
            if idx >= 0:
                chosen = l
                idx0 = idx
                break
        if chosen is not None:
            break
    assert chosen is not None
    assert idx0 is not None

    layout_part = chosen["part"]

    # Minimal template spec pointing to a bad layout name but correct layout_part.
    template_spec = {
        "version": "1.0",
        "styles": {},
        "deck": {},
        "archetypes": [
            {
                "id": "title",
                "layout": "THIS NAME DOES NOT EXIST",
                "layout_part": layout_part,
                "slots": [
                    {"name": "title", "type": "text", "required": True, "placeholder_idx": idx0},
                ],
            }
        ],
    }

    deck_spec = {"slides": [{"archetype": "title", "title": "Hello"}]}

    # Place template package on disk
    templates_root = tmp_path / "templates"
    tdir = templates_root / "t"
    tdir.mkdir(parents=True)
    (tdir / "template.json").write_text(json.dumps(template_spec), encoding="utf-8")
    (tdir / "template.pptx").write_bytes(pptx_path.read_bytes())

    out = tmp_path / "out.pptx"
    rendered = render_deck(deck_spec, template_spec, template_id="t", output_path=str(out), templates_dir=str(templates_root))

    assert Path(rendered).exists()
    out_prs = Presentation(rendered)
    assert len(out_prs.slides) == 1
