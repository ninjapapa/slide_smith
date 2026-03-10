from __future__ import annotations

import base64
import json
from pathlib import Path

from pptx import Presentation

from slide_smith.editor import add_slide_to_deck, update_slide_in_deck
from slide_smith.renderer import render_deck
from slide_smith.template_loader import load_template_spec


PNG_1X1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAusB9WwS6uoAAAAASUVORK5CYII="
)



def _make_base_deck(tmp_path: Path) -> Path:
    image_path = tmp_path / "demo.png"
    image_path.write_bytes(PNG_1X1)
    deck_spec = {
        "slides": [
            {"archetype": "title", "title": "Demo", "subtitle": "Sub"},
            {"archetype": "title_and_bullets", "title": "Highlights", "bullets": ["One", "Two"]},
        ]
    }
    template_spec = load_template_spec("default")
    output = tmp_path / "base.pptx"
    render_deck(deck_spec, template_spec, "default", str(output), base_dir=str(tmp_path))
    return output



def test_add_slide_to_deck_appends_supported_slide(tmp_path: Path) -> None:
    deck = _make_base_deck(tmp_path)
    slide_input = tmp_path / "slide.json"
    slide_input.write_text(json.dumps({"title": "Product", "body": "Body", "image": "demo.png"}))
    (tmp_path / "demo.png").write_bytes(PNG_1X1)

    add_slide_to_deck(str(deck), after_index=1, archetype="image_left_text_right", input_path=str(slide_input))

    prs = Presentation(str(deck))
    assert len(prs.slides) == 3
    assert prs.slides[2].shapes.title.text == "Product"



def test_update_slide_in_deck_updates_bullets(tmp_path: Path) -> None:
    deck = _make_base_deck(tmp_path)
    patch = tmp_path / "patch.json"
    patch.write_text(json.dumps({"title": "Updated Highlights", "bullets": ["Alpha", "Beta"]}))

    update_slide_in_deck(str(deck), index=1, input_path=str(patch))

    prs = Presentation(str(deck))
    assert prs.slides[1].shapes.title.text == "Updated Highlights"
    assert prs.slides[1].placeholders[1].text_frame.text == "Alpha\nBeta"
