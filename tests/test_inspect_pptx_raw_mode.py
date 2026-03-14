from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation

from slide_smith.commands.inspect_pptx import handle_inspect_pptx


def test_inspect_pptx_raw_mode_json(tmp_path: Path) -> None:
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    p = tmp_path / "ref.pptx"
    prs.save(p)

    code, out = handle_inspect_pptx(pptx=str(p), fmt="json", mode="raw")
    assert code == 0
    payload = json.loads(out)
    assert payload["mode"] == "raw"
    assert payload["layouts"], "expected layouts in raw mode"
    assert "slide_size" in payload
