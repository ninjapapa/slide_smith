from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from slide_smith.renderer import render_deck
from slide_smith.template_loader import load_template_spec


def test_render_deck_missing_placeholder_idx_falls_back_with_actionable_warning(tmp_path: Path) -> None:
    deck_spec = {
        "title": "Demo",
        "slides": [
            {"layout_id": "title", "title": "Demo", "subtitle": "Sub"},
        ],
    }

    template_spec = load_template_spec("default")

    # Break the template mapping to point at a non-existent placeholder index.
    for a in template_spec.get("archetypes", []):
        if a.get("id") == "title":
            for s in a.get("slots", []):
                if s.get("name") == "title":
                    s["placeholder_idx"] = 999

    render_deck(deck_spec, template_spec, "default", str(tmp_path / "out.pptx"), base_dir=str(tmp_path))

    warnings = deck_spec.get("render_warnings") or []
    assert warnings
    msg = str(warnings[0]["reason"])
    assert "idx=999" in msg
    assert "layout_id=title" in msg

    prs = Presentation(str(tmp_path / "out.pptx"))
    assert len(prs.slides) == 1
    assert prs.slides[0].shapes.title.text == "Demo"
