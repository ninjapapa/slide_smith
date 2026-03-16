from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation

from slide_smith.commands.edit_ops import handle_add_slide
from slide_smith.renderer import render_deck
from slide_smith.template_loader import load_template_spec


def _base_deck(tmp_path: Path) -> Path:
    out = tmp_path / "base.pptx"
    render_deck(
        {"slides": [{"layout_id": "title", "title": "Demo", "subtitle": "Sub"}]},
        load_template_spec("default"),
        "default",
        str(out),
        base_dir=str(tmp_path),
    )
    return out


def test_insert_slide_accepts_layout_id(tmp_path: Path) -> None:
    deck = _base_deck(tmp_path)
    slide_input = tmp_path / "slide.json"
    slide_input.write_text(json.dumps({"title": "Highlights", "bullets": ["One", "Two"]}))

    code, out = handle_add_slide(deck=str(deck), after=0, layout_id="title_and_bullets", input_path=str(slide_input))
    assert code == 0
    payload = json.loads(out)
    assert payload["status"] == "slide inserted"

    prs = Presentation(str(deck))
    assert len(prs.slides) == 2
    assert prs.slides[1].shapes.title.text == "Highlights"


def test_insert_slide_falls_back_when_layout_is_not_supported_by_template(tmp_path: Path) -> None:
    deck = _base_deck(tmp_path)
    slide_input = tmp_path / "slide.json"
    slide_input.write_text(json.dumps({"title": "Compare", "col1_body": "L", "col2_body": "R"}))

    code, out = handle_add_slide(deck=str(deck), after=0, layout_id="two_col", input_path=str(slide_input))
    assert code == 0
    payload = json.loads(out)
    assert payload["status"] == "slide inserted"

    prs = Presentation(str(deck))
    assert len(prs.slides) == 2
    assert prs.slides[1].shapes.title.text == "Compare"
    assert "L" in prs.slides[1].placeholders[1].text
