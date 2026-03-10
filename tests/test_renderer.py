from __future__ import annotations

import base64
from pathlib import Path

from pptx import Presentation

from slide_smith.renderer import render_deck
from slide_smith.styling import load_styles
from slide_smith.template_loader import load_template_spec


PNG_1X1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAusB9WwS6uoAAAAASUVORK5CYII="
)



def test_render_deck_writes_pptx(tmp_path: Path) -> None:
    image_path = tmp_path / "demo.png"
    image_path.write_bytes(PNG_1X1)

    deck_spec = {
        "title": "Demo",
        "slides": [
            {"archetype": "title", "title": "Demo", "subtitle": "Sub", "notes": "intro notes"},
            {"archetype": "title_and_bullets", "title": "Highlights", "bullets": ["One", "Two"], "notes": "bullet notes"},
            {
                "archetype": "image_left_text_right",
                "title": "Product",
                "body": "Body text",
                "image": "demo.png",
                "notes": "image notes",
            },
        ],
    }
    template_spec = load_template_spec("default")
    output = tmp_path / "out.pptx"

    render_deck(deck_spec, template_spec, "default", str(output), base_dir=str(tmp_path))

    assert output.exists()
    prs = Presentation(str(output))
    assert len(prs.slides) == 3
    assert prs.slides[0].shapes.title.text == "Demo"
    assert prs.slides[1].shapes.title.text == "Highlights"

    # Notes preserved.
    assert prs.slides[0].notes_slide.notes_text_frame.text == "intro notes"
    assert prs.slides[1].notes_slide.notes_text_frame.text == "bullet notes"
    assert prs.slides[2].notes_slide.notes_text_frame.text == "image notes"

    # Styles applied (at least for the title placeholder run).
    styles = load_styles(template_spec)
    title_run = prs.slides[0].shapes.title.text_frame.paragraphs[0].runs[0]
    assert title_run.font.name == styles["title"].font
    assert int(title_run.font.size.pt) == int(styles["title"].size_pt)
