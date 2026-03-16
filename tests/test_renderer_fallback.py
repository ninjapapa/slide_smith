from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from slide_smith.renderer import render_deck


def test_render_deck_falls_back_when_template_does_not_support_requested_layout(tmp_path: Path) -> None:
    deck_spec = {
        "title": "Demo",
        "slides": [
            {"layout_id": "two_col", "title": "Two", "col1_body": "L", "col2_body": "R"},
        ],
    }

    template_spec = {
        "template_id": "t",
        "archetypes": [
            {
                "id": "title_and_bullets",
                "layout": "Title and Content",
                "slots": [
                    {"name": "title", "type": "text", "placeholder_idx": 0},
                    {"name": "bullets", "type": "bullet_list", "placeholder_idx": 1},
                ],
            }
        ],
        "styles": {},
    }

    output = tmp_path / "out.pptx"
    render_deck(deck_spec, template_spec, "default", str(output), base_dir=str(tmp_path))

    assert output.exists()
    assert isinstance(deck_spec.get("render_warnings"), list)
    assert deck_spec["render_warnings"][0]["requested_layout_id"] == "two_col"
    assert deck_spec["render_warnings"][0]["fallback_layout_id"] == "title_and_bullets"

    prs = Presentation(str(output))
    assert len(prs.slides) == 1
    assert prs.slides[0].shapes.title.text == "Two"
    assert "L" in prs.slides[0].placeholders[1].text
    assert "R" in prs.slides[0].placeholders[1].text
