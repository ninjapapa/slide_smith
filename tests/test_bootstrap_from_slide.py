from __future__ import annotations

import base64
from pathlib import Path

from pptx import Presentation

from slide_smith.slide_instance_bootstrapper import bootstrap_archetype_from_slide


PNG_1X1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO7Wk0kAAAAASUVORK5CYII="
)


def test_bootstrap_image_left_text_right_from_slide(tmp_path: Path) -> None:
    pptx = tmp_path / "in.pptx"
    img = tmp_path / "img.png"
    img.write_bytes(PNG_1X1)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title at top-left
    slide.shapes.add_textbox(int(prs.slide_width * 0.05), int(prs.slide_height * 0.05), int(prs.slide_width * 0.6), int(prs.slide_height * 0.1)).text_frame.text = "T"
    # Body below
    slide.shapes.add_textbox(int(prs.slide_width * 0.05), int(prs.slide_height * 0.25), int(prs.slide_width * 0.5), int(prs.slide_height * 0.5)).text_frame.text = "B"
    # Image on right
    slide.shapes.add_picture(str(img), int(prs.slide_width * 0.6), int(prs.slide_height * 0.2), width=int(prs.slide_width * 0.35), height=int(prs.slide_height * 0.6))

    prs.save(str(pptx))

    boot = bootstrap_archetype_from_slide(str(pptx), slide_number=1, archetype_id="image_left_text_right")
    spec = boot.archetype_spec

    assert spec["id"] == "image_left_text_right"
    slots = {s["name"]: s for s in spec["slots"]}
    assert "box" in slots["title"]
    assert "box" in slots["body"]
    assert "box" in slots["image"]

    # Basic sanity: image box should be on right half.
    assert slots["image"]["box"]["x"] > 0.5
