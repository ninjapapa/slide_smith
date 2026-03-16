from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation

from slide_smith.deck_spec import ARCHETYPE_ALIASES, normalize_deck_spec
from slide_smith.styling import apply_text_style, load_styles
from slide_smith.template_loader import template_dir


class RenderingError(Exception):
    """Raised when a deck cannot be rendered."""


FALLBACK_LAYOUT_ID = "title_and_bullets"


def _presentation_for_template(
    template_id: str,
    templates_dir: str | None = None,
    *,
    preserve_template_slides: bool = False,
) -> Presentation:
    """Load a Presentation for rendering.

    By default we treat template.pptx as a *theme/layout source* and start with zero
    content slides (rich templates often include many sample slides that should not
    be preserved in generated outputs).

    If preserve_template_slides=True, keep any existing slides from template.pptx.
    """

    pptx_path = template_dir(template_id, templates_dir) / "template.pptx"
    if not pptx_path.exists():
        return Presentation()

    prs = Presentation(str(pptx_path))

    if preserve_template_slides:
        return prs

    # Drop all existing slides (keep masters/layouts/theme).
    try:
        from slide_smith.pptx_edit import delete_slide

        while len(prs.slides) > 0:
            delete_slide(prs, 0)
    except Exception as exc:
        raise RenderingError(f"Failed to clear template slides: {exc}") from exc

    return prs


def _layout_by_name(prs: Presentation, name: str):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    raise RenderingError(f"Slide layout '{name}' not found in template presentation")


def _fallback_text_lines(slide_spec: dict[str, Any]) -> list[str]:
    lines: list[str] = []

    def add(value: object) -> None:
        if isinstance(value, str) and value.strip():
            lines.append(value.strip())

    subtitle = slide_spec.get("subtitle")
    add(subtitle)

    body = slide_spec.get("body")
    add(body)

    for i in range(1, 6):
        add(slide_spec.get(f"col{i}_title"))
        add(slide_spec.get(f"col{i}_body"))
        add(slide_spec.get(f"pillar{i}_body"))
        add(slide_spec.get(f"item{i}_body"))

    bullets = slide_spec.get("bullets")
    if isinstance(bullets, list):
        for item in bullets:
            add(item)

    for key in ("left", "right"):
        side = slide_spec.get(key)
        if isinstance(side, dict):
            add(side.get("title"))
            add(side.get("body"))

    items = slide_spec.get("items")
    if isinstance(items, list):
        for item in items:
            if isinstance(item, dict):
                for key in ("marker", "title", "heading", "body", "caption", "label"):
                    add(item.get(key))

    for key in ("left_image", "right_image"):
        if key in slide_spec:
            add(f"[{key}: provided]")

    if not lines:
        lines.append("Content could not be rendered in the requested layout.")

    return lines


def _make_fallback_slide_spec(slide_spec: dict[str, Any], *, requested_layout_id: str, reason: str) -> dict[str, Any]:
    title = slide_spec.get("title")
    if not isinstance(title, str) or not title.strip():
        title = f"Fallback for {requested_layout_id}"

    bullets = _fallback_text_lines(slide_spec)
    out = dict(slide_spec)
    out["archetype"] = FALLBACK_LAYOUT_ID
    out["layout_id"] = FALLBACK_LAYOUT_ID
    out["title"] = title
    out["bullets"] = bullets
    out.pop("image", None)
    return out


def _record_render_warning(deck_spec: dict[str, Any], *, slide_index: int, requested_layout_id: str, reason: str) -> None:
    warnings = deck_spec.setdefault("render_warnings", [])
    if not isinstance(warnings, list):
        return
    warnings.append(
        {
            "slide_index": slide_index,
            "requested_layout_id": requested_layout_id,
            "fallback_layout_id": FALLBACK_LAYOUT_ID,
            "reason": reason,
            "status": "fallback",
        }
    )


def _layout_for_archetype(prs: Presentation, archetype_spec: dict[str, Any]):
    """Resolve a layout for an archetype.

    Prefer `layout_part` when available (more stable for rich branded templates).
    Fallback to `layout` name.
    """

    from slide_smith.layout_resolver import LayoutResolveError, resolve_layout

    layout_name = archetype_spec.get("layout")
    layout_part = archetype_spec.get("layout_part")

    try:
        return resolve_layout(prs=prs, layout_name=layout_name, layout_part=layout_part).layout
    except LayoutResolveError as exc:
        raise RenderingError(str(exc)) from exc


def _slot_spec(archetype_spec: dict[str, Any], slot_name: str) -> dict[str, Any] | None:
    for slot in archetype_spec.get("slots", []):
        if isinstance(slot, dict) and slot.get("name") == slot_name:
            return slot
    return None


def _slot_index(archetype_spec: dict[str, Any], slot_name: str, default: int | None = None) -> int | None:
    """Resolve placeholder_idx for a slot.

    `default` is only used when the slot is missing entirely.
    If the slot exists but does not declare placeholder_idx (e.g. box-based slots),
    this returns None.
    """

    slot = _slot_spec(archetype_spec, slot_name)
    if slot is None:
        return default
    idx = slot.get("placeholder_idx")
    if isinstance(idx, int):
        return idx
    return None


def _slot_box(archetype_spec: dict[str, Any], slot_name: str) -> dict[str, Any] | None:
    slot = _slot_spec(archetype_spec, slot_name)
    if slot is None:
        return None
    box = slot.get("box")
    if isinstance(box, dict):
        return box
    return None


def _box_to_emu(box: dict[str, Any], *, slide_w_emu: int, slide_h_emu: int) -> tuple[int, int, int, int] | None:
    """Convert a slot box to EMU coordinates.

    Supported forms:
    - units=relative: x,y,w,h are floats in [0,1]
    - units=emu: left,top,width,height are ints (or x,y,w,h)

    Returns: (left, top, width, height)
    """

    units = box.get("units")
    if units == "relative":
        try:
            x = float(box["x"])
            y = float(box["y"])
            w = float(box["w"])
            h = float(box["h"])
        except Exception:
            return None
        left = int(x * slide_w_emu)
        top = int(y * slide_h_emu)
        width = int(w * slide_w_emu)
        height = int(h * slide_h_emu)
        return left, top, width, height

    if units == "emu":
        # Accept either (left,top,width,height) or (x,y,w,h)
        def gi(*keys):
            for k in keys:
                v = box.get(k)
                if isinstance(v, int):
                    return v
            return None

        left = gi("left", "x")
        top = gi("top", "y")
        width = gi("width", "w")
        height = gi("height", "h")
        if None in (left, top, width, height):
            return None
        return int(left), int(top), int(width), int(height)

    return None


def _required_slot_target(
    archetype_id: str,
    archetype_spec: dict[str, Any],
    slot_name: str,
    *,
    default_idx: int | None = None,
) -> tuple[int | None, dict[str, Any] | None]:
    slot = _slot_spec(archetype_spec, slot_name)
    idx = _slot_index(archetype_spec, slot_name, default_idx)
    box = _slot_box(archetype_spec, slot_name)

    required = bool(slot.get("required")) if slot is not None else False
    if required and idx is None and box is None:
        raise RenderingError(
            f"Template archetype '{archetype_id}' missing required slot mapping '{slot_name}' (placeholder_idx or box)"
        )

    return idx, box


def _set_placeholder_text(slide, idx: int | None, text: str | None, style=None, *, context: str = "") -> None:
    if idx is None or text is None:
        return
    try:
        ph = slide.placeholders[idx]
        ph.text = text
        if style is not None and hasattr(ph, "text_frame"):
            apply_text_style(ph.text_frame, style)
    except KeyError as exc:
        ctx = f" ({context})" if context else ""
        raise RenderingError(f"Placeholder idx={idx} not found on slide{ctx}") from exc


def _set_box_text(slide, box_emu: tuple[int, int, int, int] | None, text: str | None, style=None) -> None:
    if box_emu is None or text is None:
        return
    left, top, width, height = box_emu
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.text = text
    if style is not None:
        apply_text_style(tf, style)


def _set_box_bullets(slide, box_emu: tuple[int, int, int, int] | None, bullets: list[str], style=None) -> None:
    if box_emu is None:
        return
    left, top, width, height = box_emu
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    if not bullets:
        tf.text = ""
        return
    tf.text = bullets[0]
    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
    if style is not None:
        apply_text_style(tf, style)


def _set_box_image(slide, box_emu: tuple[int, int, int, int] | None, image_path: Path | None) -> None:
    if box_emu is None or image_path is None:
        return
    left, top, width, height = box_emu
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def _set_slot_text(
    slide,
    archetype_id: str,
    archetype_spec: dict[str, Any],
    slot_name: str,
    text: str | None,
    style,
    *,
    slide_w_emu: int,
    slide_h_emu: int,
    default_idx: int | None = None,
    context: str = "",
) -> None:
    idx, box = _required_slot_target(archetype_id, archetype_spec, slot_name, default_idx=default_idx)
    if idx is not None:
        _set_placeholder_text(slide, idx, text, style, context=context)
        return
    if box is not None:
        box_emu = _box_to_emu(box, slide_w_emu=slide_w_emu, slide_h_emu=slide_h_emu)
        _set_box_text(slide, box_emu, text, style)


def _resolve_image_path(base_dir: Path, image_field: Any) -> Path | None:
    image_path = None
    if isinstance(image_field, str):
        image_path = image_field
    elif isinstance(image_field, dict):
        image_path = image_field.get("path")

    if not image_path:
        return None

    p = Path(str(image_path)).expanduser()
    resolved = (base_dir / p).resolve() if not p.is_absolute() else p.resolve()
    if not resolved.exists():
        raise RenderingError(f"Image file not found: {resolved}")
    if not resolved.is_file():
        raise RenderingError(f"Image path is not a file: {resolved}")
    return resolved


def _set_slot_image(
    slide,
    archetype_id: str,
    archetype_spec: dict[str, Any],
    slot_name: str,
    image_path: Path | None,
    *,
    slide_w_emu: int,
    slide_h_emu: int,
    default_idx: int | None = None,
    context: str = "",
) -> None:
    if image_path is None:
        return

    idx, box = _required_slot_target(archetype_id, archetype_spec, slot_name, default_idx=default_idx)
    if idx is not None:
        try:
            image_placeholder = slide.placeholders[idx]
        except KeyError as exc:
            ctx = f" ({context})" if context else ""
            raise RenderingError(f"Placeholder idx={idx} not found on slide{ctx}") from exc
        slide.shapes.add_picture(
            str(image_path),
            image_placeholder.left,
            image_placeholder.top,
            width=image_placeholder.width,
            height=image_placeholder.height,
        )
        return

    if box is not None:
        box_emu = _box_to_emu(box, slide_w_emu=slide_w_emu, slide_h_emu=slide_h_emu)
        _set_box_image(slide, box_emu, image_path)


def _render_title(slide, spec: dict[str, Any], styles, archetype_spec: dict[str, Any], archetype_id: str, *, slide_w_emu: int, slide_h_emu: int) -> None:
    _set_slot_text(
        slide,
        archetype_id,
        archetype_spec,
        "title",
        spec.get("title", ""),
        styles.get("title"),
        slide_w_emu=slide_w_emu,
        slide_h_emu=slide_h_emu,
        default_idx=0,
        context=f"archetype={archetype_id} slot=title",
    )
    _set_slot_text(
        slide,
        archetype_id,
        archetype_spec,
        "subtitle",
        spec.get("subtitle", ""),
        styles.get("subtitle"),
        slide_w_emu=slide_w_emu,
        slide_h_emu=slide_h_emu,
        default_idx=1,
        context=f"archetype={archetype_id} slot=subtitle",
    )


def _render_section(slide, spec: dict[str, Any], styles, archetype_spec: dict[str, Any], archetype_id: str, *, slide_w_emu: int, slide_h_emu: int) -> None:
    _set_slot_text(
        slide,
        archetype_id,
        archetype_spec,
        "title",
        spec.get("title", ""),
        styles.get("title"),
        slide_w_emu=slide_w_emu,
        slide_h_emu=slide_h_emu,
        default_idx=0,
        context=f"archetype={archetype_id} slot=title",
    )

    subtitle = spec.get("subtitle") or spec.get("body", "")
    subtitle_idx = _slot_index(archetype_spec, "subtitle")
    body_idx = _slot_index(archetype_spec, "body")
    idx = subtitle_idx if subtitle_idx is not None else body_idx
    box = _slot_box(archetype_spec, "subtitle") or _slot_box(archetype_spec, "body")

    if idx is not None:
        _set_placeholder_text(
            slide,
            idx,
            subtitle,
            styles.get("subtitle"),
            context=f"archetype={archetype_id} slot=subtitle/body",
        )
        return

    if box is not None:
        box_emu = _box_to_emu(box, slide_w_emu=slide_w_emu, slide_h_emu=slide_h_emu)
        _set_box_text(slide, box_emu, subtitle, styles.get("subtitle"))
        return

    # If neither placeholder nor box is provided, keep behavior consistent: no-op.


def _render_title_and_bullets(slide, spec: dict[str, Any], styles, archetype_spec: dict[str, Any], archetype_id: str, *, slide_w_emu: int, slide_h_emu: int) -> None:
    _set_slot_text(
        slide,
        archetype_id,
        archetype_spec,
        "title",
        spec.get("title", ""),
        styles.get("title"),
        slide_w_emu=slide_w_emu,
        slide_h_emu=slide_h_emu,
        default_idx=0,
        context=f"archetype={archetype_id} slot=title",
    )

    # bullets/body target
    body_idx = _slot_index(archetype_spec, "bullets")
    if body_idx is None:
        body_idx = _slot_index(archetype_spec, "body", 1)

    body_box = _slot_box(archetype_spec, "bullets") or _slot_box(archetype_spec, "body")

    bullets = spec.get("bullets") or []
    if not isinstance(bullets, list):
        bullets = []

    if body_idx is None and body_box is None:
        # If the template explicitly says bullets is required, fail loudly.
        _required_slot_target(archetype_id, archetype_spec, "bullets")
        raise RenderingError(
            f"Cannot resolve placeholder/box for slot 'bullets/body' in archetype '{archetype_id}'"
        )

    if body_idx is not None:
        try:
            body_placeholder = slide.placeholders[body_idx]
        except KeyError as exc:
            raise RenderingError(
                f"Placeholder idx={body_idx} not found on slide (archetype={archetype_id} slot=bullets/body)"
            ) from exc
        text_frame = body_placeholder.text_frame
        if not bullets:
            text_frame.text = spec.get("body", "")
            apply_text_style(text_frame, styles.get("body"))
            return
        text_frame.text = bullets[0]
        for bullet in bullets[1:]:
            paragraph = text_frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
        apply_text_style(text_frame, styles.get("bullets"))
        return

    # box-based bullets
    box_emu = _box_to_emu(body_box, slide_w_emu=slide_w_emu, slide_h_emu=slide_h_emu) if body_box else None
    if not bullets:
        _set_box_text(slide, box_emu, spec.get("body", ""), styles.get("body"))
    else:
        _set_box_bullets(slide, box_emu, bullets, styles.get("bullets"))


def _render_image_left_text_right(
    slide,
    spec: dict[str, Any],
    base_dir: Path,
    styles,
    archetype_spec: dict[str, Any],
    archetype_id: str,
    *,
    slide_w_emu: int,
    slide_h_emu: int,
) -> None:
    _set_slot_text(
        slide,
        archetype_id,
        archetype_spec,
        "title",
        spec.get("title", ""),
        styles.get("title"),
        slide_w_emu=slide_w_emu,
        slide_h_emu=slide_h_emu,
        default_idx=0,
        context=f"archetype={archetype_id} slot=title",
    )

    resolved_image = _resolve_image_path(base_dir, spec.get("image"))

    _set_slot_image(
        slide,
        archetype_id,
        archetype_spec,
        "image",
        resolved_image,
        slide_w_emu=slide_w_emu,
        slide_h_emu=slide_h_emu,
        default_idx=1,
        context=f"archetype={archetype_id} slot=image",
    )

    body_idx, body_box = _required_slot_target(archetype_id, archetype_spec, "body", default_idx=2)

    # body
    body_text = spec.get("body", "")
    if body_idx is not None:
        try:
            body_placeholder = slide.placeholders[body_idx]
        except KeyError as exc:
            raise RenderingError(
                f"Placeholder idx={body_idx} not found on slide (archetype={archetype_id} slot=body)"
            ) from exc
        body_placeholder.text = body_text
        apply_text_style(body_placeholder.text_frame, styles.get("body"))
    elif body_box is not None:
        box_emu = _box_to_emu(body_box, slide_w_emu=slide_w_emu, slide_h_emu=slide_h_emu)
        _set_box_text(slide, box_emu, body_text, styles.get("body"))


def _render_v2_families(
    slide,
    spec: dict[str, Any],
    base_dir: Path,
    styles,
    archetype_spec: dict[str, Any],
    archetype_id: str,
    *,
    slide_w_emu: int,
    slide_h_emu: int,
) -> None:
    """MVP renderer for proposed v2 archetype families.

    Notes:
    - This is intentionally conservative/minimal.
    - These archetypes are not yet part of the v1 deck-spec schema; they are introduced
      for forward compatibility + early template experiments.
    - Slot mapping is still driven entirely by template.json.

    Expected template slots (convention):
    - message: title, body (or quote), attribution (optional)
    - multi_col: title, col1_body..col4_body (best-effort)
    - image_text: title, image, body (image_side is a param only; mapping decides actual layout)
    - list_visual: title, bullets (preferred) OR body (fallback)
    - metrics: title, body (metrics rendered as text for MVP)
    """

    def set_text_slot(slot_name: str, value: str | None, style_key: str = "body"):
        _set_slot_text(
            slide,
            archetype_id,
            archetype_spec,
            slot_name,
            value,
            styles.get(style_key),
            slide_w_emu=slide_w_emu,
            slide_h_emu=slide_h_emu,
            context=f"archetype={archetype_id} slot={slot_name}",
        )

    # Title is shared across the families.
    _set_slot_text(
        slide,
        archetype_id,
        archetype_spec,
        "title",
        spec.get("title", ""),
        styles.get("title"),
        slide_w_emu=slide_w_emu,
        slide_h_emu=slide_h_emu,
        default_idx=0,
        context=f"archetype={archetype_id} slot=title",
    )

    if archetype_id == "message":
        # Prefer explicit body/quote; templates can choose which slot they expose.
        body = spec.get("body") or spec.get("quote") or ""
        attribution = spec.get("attribution") or ""

        # Try to populate both body and quote if template declares them.
        if _slot_spec(archetype_spec, "body") is not None:
            set_text_slot("body", body, style_key="body")
        if _slot_spec(archetype_spec, "quote") is not None:
            set_text_slot("quote", body, style_key="body")
        if _slot_spec(archetype_spec, "attribution") is not None:
            set_text_slot("attribution", attribution, style_key="subtitle")
        elif _slot_spec(archetype_spec, "subtitle") is not None:
            set_text_slot("subtitle", attribution, style_key="subtitle")
        return

    if archetype_id == "multi_col":
        # Canonical v2 shape: items[] of objects with body/heading/label.
        # MVP: render each item as text in colN_body slots.
        items = spec.get("items") or []
        if not isinstance(items, list):
            items = []
        # Fallback to legacy colN_body if user provided it.
        if not items:
            for i in range(1, 5):
                k = f"col{i}_body"
                if k in spec and isinstance(spec.get(k), str):
                    items.append({"body": spec.get(k)})

        for i, item in enumerate(items[:4], start=1):
            if not isinstance(item, dict):
                continue
            heading = item.get("heading")
            body = item.get("body")
            label = item.get("label") or item.get("number")
            parts = []
            if isinstance(label, str) and label:
                parts.append(label)
            if isinstance(heading, str) and heading:
                parts.append(heading)
            if isinstance(body, str) and body:
                parts.append(body)
            text = "\n".join(parts)
            set_text_slot(f"col{i}_body", text, style_key="body")
        return

    if archetype_id == "image_text":
        # Same semantics as image_left_text_right, but param-driven and template decides layout.
        # MVP implementation: reuse image rendering logic; expect slots: image + body.
        image_field = spec.get("image")
        image_path = None
        if isinstance(image_field, str):
            image_path = image_field
        elif isinstance(image_field, dict):
            image_path = image_field.get("path")

        image_idx, image_box = _required_slot_target(archetype_id, archetype_spec, "image", default_idx=1)
        body_idx, body_box = _required_slot_target(archetype_id, archetype_spec, "body", default_idx=2)

        resolved_image: Path | None = None
        if image_path:
            p = Path(str(image_path)).expanduser()
            resolved_image = (base_dir / p).resolve() if not p.is_absolute() else p.resolve()
            if not resolved_image.exists():
                raise RenderingError(f"Image file not found: {resolved_image}")
            if not resolved_image.is_file():
                raise RenderingError(f"Image path is not a file: {resolved_image}")

        if resolved_image is not None:
            if image_idx is not None:
                try:
                    image_placeholder = slide.placeholders[image_idx]
                except KeyError as exc:
                    raise RenderingError(
                        f"Placeholder idx={image_idx} not found on slide (archetype={archetype_id} slot=image)"
                    ) from exc
                slide.shapes.add_picture(
                    str(resolved_image),
                    image_placeholder.left,
                    image_placeholder.top,
                    width=image_placeholder.width,
                    height=image_placeholder.height,
                )
            elif image_box is not None:
                box_emu = _box_to_emu(image_box, slide_w_emu=slide_w_emu, slide_h_emu=slide_h_emu)
                _set_box_image(slide, box_emu, resolved_image)

        body_text = spec.get("body", "")
        if body_idx is not None:
            try:
                body_placeholder = slide.placeholders[body_idx]
            except KeyError as exc:
                raise RenderingError(
                    f"Placeholder idx={body_idx} not found on slide (archetype={archetype_id} slot=body)"
                ) from exc
            body_placeholder.text = body_text
            apply_text_style(body_placeholder.text_frame, styles.get("body"))
        elif body_box is not None:
            box_emu = _box_to_emu(body_box, slide_w_emu=slide_w_emu, slide_h_emu=slide_h_emu)
            _set_box_text(slide, box_emu, body_text, styles.get("body"))
        return

    if archetype_id == "list_visual":
        # MVP: treat as bullets/body.
        items = spec.get("items") or []
        bullets = []
        if isinstance(items, list):
            for it in items:
                if not isinstance(it, dict):
                    continue
                label = it.get("label") or it.get("number")
                body = it.get("body")
                if isinstance(label, str) and label and isinstance(body, str) and body:
                    bullets.append(f"{label} {body}")
                elif isinstance(body, str) and body:
                    bullets.append(body)

        # Prefer a bullets slot if present; otherwise fall back to body.
        if _slot_spec(archetype_spec, "bullets") is not None:
            # Reuse the bullets helper by delegating to title_and_bullets-ish behavior.
            proxy = dict(spec)
            proxy["bullets"] = bullets
            _render_title_and_bullets(slide, proxy, styles, archetype_spec, archetype_id, slide_w_emu=slide_w_emu, slide_h_emu=slide_h_emu)
        else:
            set_text_slot("body", "\n".join(bullets), style_key="body")
        return

    if archetype_id == "metrics":
        ms = spec.get("metrics") or []
        lines = []
        if isinstance(ms, list):
            for m in ms:
                if not isinstance(m, dict):
                    continue
                value = m.get("value")
                label = m.get("label")
                detail = m.get("detail")
                if isinstance(value, str) and isinstance(label, str):
                    s = f"{value} — {label}" if label else value
                    if isinstance(detail, str) and detail:
                        s += f"\n{detail}"
                    lines.append(s)
        set_text_slot("body", "\n\n".join(lines), style_key="body")
        return

    raise RenderingError(f"v2 archetype family '{archetype_id}' is not implemented")


def _render_extended(
    slide,
    spec: dict[str, Any],
    base_dir: Path,
    styles,
    archetype_spec: dict[str, Any],
    archetype_id: str,
    *,
    slide_w_emu: int,
    slide_h_emu: int,
) -> None:
    """MVP renderer for v1.1 extended archetypes.

    This implementation is intentionally placeholder-driven and template-mapping-driven.
    It treats the template spec slot mappings as the source of truth.

    v1.2+ note: box-based slots are supported for text fields as a fallback.
    """

    title = spec.get("title", "")
    _set_slot_text(
        slide,
        archetype_id,
        archetype_spec,
        "title",
        title,
        styles.get("title"),
        slide_w_emu=slide_w_emu,
        slide_h_emu=slide_h_emu,
        default_idx=0,
        context=f"archetype={archetype_id} slot=title",
    )

    def set_text_slot(slot_name: str, value: str | None, style_key: str = "body"):
        _set_slot_text(
            slide,
            archetype_id,
            archetype_spec,
            slot_name,
            value,
            styles.get(style_key),
            slide_w_emu=slide_w_emu,
            slide_h_emu=slide_h_emu,
            context=f"archetype={archetype_id} slot={slot_name}",
        )

    if archetype_id in {"two_col", "three_col", "four_col"}:
        n = {"two_col": 2, "three_col": 3, "four_col": 4}[archetype_id]
        for i in range(1, n + 1):
            set_text_slot(f"col{i}_body", spec.get(f"col{i}_body"))
        return

    if archetype_id == "two_col_with_subtitle":
        set_text_slot("subtitle", spec.get("subtitle"), style_key="subtitle")
        for i in (1, 2):
            set_text_slot(f"col{i}_body", spec.get(f"col{i}_body"))
        return

    if archetype_id in {"pillars_3", "pillars_4"}:
        n = 3 if archetype_id == "pillars_3" else 4
        for i in range(1, n + 1):
            set_text_slot(f"pillar{i}_body", spec.get(f"pillar{i}_body"))
        return

    if archetype_id in {"table", "version_page"}:
        set_text_slot("table_text", spec.get("table_text"), style_key="body")
        return

    if archetype_id == "table_plus_description":
        set_text_slot("table_text", spec.get("table_text"), style_key="body")
        set_text_slot("body", spec.get("body"), style_key="body")
        return

    if archetype_id == "timeline_horizontal":
        for i in range(1, 11):
            k = f"milestone{i}_body"
            if k in spec:
                set_text_slot(k, spec.get(k))
        return

    if archetype_id == "title_subtitle":
        set_text_slot("subtitle", spec.get("subtitle"), style_key="subtitle")
        return

    if archetype_id == "title_only_freeform":
        return

    if archetype_id == "picture_compare":
        left = spec.get("left") or {}
        right = spec.get("right") or {}

        if isinstance(left, dict):
            _set_slot_image(
                slide,
                archetype_id,
                archetype_spec,
                "left_image",
                _resolve_image_path(base_dir, left.get("image")),
                slide_w_emu=slide_w_emu,
                slide_h_emu=slide_h_emu,
                context=f"archetype={archetype_id} slot=left_image",
            )
            set_text_slot("left_title", left.get("title"), style_key="body")
            set_text_slot("left_body", left.get("body"), style_key="body")

        if isinstance(right, dict):
            _set_slot_image(
                slide,
                archetype_id,
                archetype_spec,
                "right_image",
                _resolve_image_path(base_dir, right.get("image")),
                slide_w_emu=slide_w_emu,
                slide_h_emu=slide_h_emu,
                context=f"archetype={archetype_id} slot=right_image",
            )
            set_text_slot("right_title", right.get("title"), style_key="body")
            set_text_slot("right_body", right.get("body"), style_key="body")
        return

    if archetype_id == "agenda_with_image":
        _set_slot_image(
            slide,
            archetype_id,
            archetype_spec,
            "image",
            _resolve_image_path(base_dir, spec.get("image")),
            slide_w_emu=slide_w_emu,
            slide_h_emu=slide_h_emu,
            context=f"archetype={archetype_id} slot=image",
        )

        items = spec.get("items") or []

        # Preferred: map into dedicated placeholders when the template exposes them.
        # Convention:
        # - item{n}_marker (optional)
        # - item{n}_body (required)
        # If the template doesn't have these slots, we fall back to bullets/body.
        has_item_slots = _slot_spec(archetype_spec, "item1_body") is not None
        if has_item_slots and isinstance(items, list):
            for idx, it in enumerate(items, start=1):
                if not isinstance(it, dict):
                    continue
                marker = it.get("marker")
                body = it.get("body")

                # If marker is missing, auto-number.
                if not (isinstance(marker, str) and marker):
                    marker = str(idx)

                if isinstance(body, str) and body:
                    # marker slot is optional; only set if it exists.
                    if _slot_spec(archetype_spec, f"item{idx}_marker") is not None:
                        set_text_slot(f"item{idx}_marker", str(marker), style_key="body")
                    set_text_slot(f"item{idx}_body", body, style_key="body")
            return

        # Fallback: render as bullets/body.
        lines: list[str] = []
        if isinstance(items, list):
            for it in items:
                if not isinstance(it, dict):
                    continue
                marker = it.get("marker")
                body = it.get("body")
                if isinstance(body, str) and body:
                    if isinstance(marker, str) and marker:
                        lines.append(f"{marker} {body}")
                    else:
                        lines.append(body)

        # Prefer bullets slot when available; otherwise fall back to body.
        if _slot_spec(archetype_spec, "bullets") is not None:
            proxy = dict(spec)
            proxy["bullets"] = lines
            _render_title_and_bullets(slide, proxy, styles, archetype_spec, archetype_id, slide_w_emu=slide_w_emu, slide_h_emu=slide_h_emu)
        else:
            set_text_slot("body", "\n".join(lines), style_key="body")
        return

    if archetype_id in {"three_col_with_subtitle", "three_col_with_icons", "five_col_with_icons"}:
        items = spec.get("items") or []
        if archetype_id == "three_col_with_subtitle":
            set_text_slot("subtitle", spec.get("subtitle"), style_key="subtitle")

        if not isinstance(items, list):
            items = []

        # Convention-based mapping into slot names.
        for idx, it in enumerate(items, start=1):
            if not isinstance(it, dict):
                continue

            if archetype_id == "three_col_with_subtitle":
                set_text_slot(f"col{idx}_title", it.get("title"), style_key="body")
                set_text_slot(f"col{idx}_body", it.get("body"), style_key="body")
            elif archetype_id == "three_col_with_icons":
                set_text_slot(f"col{idx}_title", it.get("title"), style_key="body")
                set_text_slot(f"col{idx}_body", it.get("body"), style_key="body")
                set_text_slot(f"col{idx}_caption", it.get("caption"), style_key="body")
                _set_slot_image(
                    slide,
                    archetype_id,
                    archetype_spec,
                    f"col{idx}_icon",
                    _resolve_image_path(base_dir, it.get("icon")),
                    slide_w_emu=slide_w_emu,
                    slide_h_emu=slide_h_emu,
                    context=f"archetype={archetype_id} slot=col{idx}_icon",
                )
            else:  # five_col_with_icons
                set_text_slot(f"item{idx}_body", it.get("body"), style_key="body")
                _set_slot_image(
                    slide,
                    archetype_id,
                    archetype_spec,
                    f"item{idx}_icon",
                    _resolve_image_path(base_dir, it.get("icon")),
                    slide_w_emu=slide_w_emu,
                    slide_h_emu=slide_h_emu,
                    context=f"archetype={archetype_id} slot=item{idx}_icon",
                )
        return

    raise RenderingError(f"Extended archetype '{archetype_id}' is not implemented")


def _set_notes(slide, notes: str | None) -> None:
    if not notes:
        return
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes


def render_deck(
    deck_spec: dict[str, Any],
    template_spec: dict[str, Any],
    template_id: str,
    output_path: str,
    base_dir: str | None = None,
    templates_dir: str | None = None,
) -> str:
    original_deck_spec = deck_spec
    deck_spec, _ = normalize_deck_spec(deck_spec)
    # Default behavior: do NOT preserve existing template slides.
    # This avoids rich templates polluting the output with sample pages.
    preserve = bool(deck_spec.get("preserve_template_slides"))
    prs = _presentation_for_template(template_id, templates_dir=templates_dir, preserve_template_slides=preserve)
    source_dir = Path(base_dir or ".").resolve()

    slide_w_emu = int(prs.slide_width)
    slide_h_emu = int(prs.slide_height)

    styles = load_styles(template_spec)
    archetypes = {item["id"]: item for item in template_spec.get("archetypes", []) if isinstance(item, dict) and isinstance(item.get("id"), str)}

    deck_meta = template_spec.get("deck") or {}
    native_pref = deck_meta.get("native_preferred") or {}

    def resolve_template_archetype_id(slide_archetype_id: str) -> str:
        """Resolve the template archetype to use for a given slide archetype.

        Allows templates to define a preferred template-native archetype for a core archetype.
        The slide archetype (and therefore renderer code path) remains the slide's archetype;
        only the *template mapping* (layout + slot placeholder indices/boxes) can be swapped.
        """

        if isinstance(native_pref, dict):
            preferred = native_pref.get(slide_archetype_id)
            if isinstance(preferred, str) and preferred in archetypes:
                return preferred

        if slide_archetype_id in archetypes:
            return slide_archetype_id

        # Compatibility: if the normalized/preferred id is missing, allow templates
        # that still define the legacy alias id.
        for legacy_id, preferred_id in ARCHETYPE_ALIASES.items():
            if slide_archetype_id == preferred_id and legacy_id in archetypes:
                return legacy_id

        return slide_archetype_id

    def render_one(slide_spec: dict[str, Any]) -> None:
        archetype = slide_spec["archetype"]
        template_archetype_id = resolve_template_archetype_id(archetype)
        if template_archetype_id not in archetypes:
            raise RenderingError(
                f"Archetype '{archetype}' not supported by template '{template_id}'"
                + (" (native_preferred mapping pointed to missing archetype)" if template_archetype_id != archetype else "")
            )
        archetype_spec = archetypes[template_archetype_id]
        slide = prs.slides.add_slide(_layout_for_archetype(prs, archetype_spec))

        if archetype == "title":
            _render_title(slide, slide_spec, styles, archetype_spec, archetype, slide_w_emu=slide_w_emu, slide_h_emu=slide_h_emu)
        elif archetype == "section":
            _render_section(slide, slide_spec, styles, archetype_spec, archetype, slide_w_emu=slide_w_emu, slide_h_emu=slide_h_emu)
        elif archetype == "title_and_bullets":
            _render_title_and_bullets(slide, slide_spec, styles, archetype_spec, archetype, slide_w_emu=slide_w_emu, slide_h_emu=slide_h_emu)
        elif archetype in {"image_left_text_right", "text_with_image"}:
            _render_image_left_text_right(slide, slide_spec, source_dir, styles, archetype_spec, archetype, slide_w_emu=slide_w_emu, slide_h_emu=slide_h_emu)
        elif archetype == "title_subtitle_and_bullets":
            _render_title_and_bullets(slide, slide_spec, styles, archetype_spec, archetype, slide_w_emu=slide_w_emu, slide_h_emu=slide_h_emu)
            _set_slot_text(
                slide,
                archetype,
                archetype_spec,
                "subtitle",
                slide_spec.get("subtitle"),
                styles.get("subtitle"),
                slide_w_emu=slide_w_emu,
                slide_h_emu=slide_h_emu,
                context=f"archetype={archetype} slot=subtitle",
            )
        elif archetype in {
            "two_col",
            "three_col",
            "four_col",
            "pillars_3",
            "pillars_4",
            "table",
            "table_plus_description",
            "timeline_horizontal",
            # redesign extended
            "title_subtitle",
            "version_page",
            "agenda_with_image",
            "two_col_with_subtitle",
            "three_col_with_subtitle",
            "three_col_with_icons",
            "five_col_with_icons",
            "picture_compare",
            "title_only_freeform",
        }:
            _render_extended(
                slide,
                slide_spec,
                source_dir,
                styles,
                archetype_spec,
                archetype,
                slide_w_emu=slide_w_emu,
                slide_h_emu=slide_h_emu,
            )
        elif archetype in {"message", "multi_col", "image_text", "list_visual", "metrics"}:
            _render_v2_families(slide, slide_spec, source_dir, styles, archetype_spec, archetype, slide_w_emu=slide_w_emu, slide_h_emu=slide_h_emu)
        else:
            raise RenderingError(f"Archetype '{archetype}' is not implemented")

        _set_notes(slide, slide_spec.get("notes"))

    for i, slide_spec in enumerate(deck_spec.get("slides", [])):
        before_count = len(prs.slides)
        try:
            render_one(slide_spec)
        except RenderingError as exc:
            # Remove any partially-added slide before fallback so output slide count stays stable.
            try:
                from slide_smith.pptx_edit import delete_slide

                while len(prs.slides) > before_count:
                    delete_slide(prs, len(prs.slides) - 1)
            except Exception:
                pass

            requested_layout_id = str(slide_spec.get("layout_id") or slide_spec.get("archetype") or "unknown")
            if requested_layout_id == FALLBACK_LAYOUT_ID:
                raise
            _record_render_warning(deck_spec, slide_index=i, requested_layout_id=requested_layout_id, reason=str(exc))
            fallback_spec = _make_fallback_slide_spec(slide_spec, requested_layout_id=requested_layout_id, reason=str(exc))
            render_one(fallback_spec)

    if isinstance(original_deck_spec, dict):
        if isinstance(deck_spec.get("render_warnings"), list):
            original_deck_spec["render_warnings"] = deck_spec.get("render_warnings")
        # Preserve normalized slides for callers that pass layout_id-only input.
        if isinstance(deck_spec.get("slides"), list):
            original_deck_spec["slides"] = deck_spec.get("slides")

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return str(out)
