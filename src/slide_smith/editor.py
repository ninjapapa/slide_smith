from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from pptx import Presentation

from slide_smith.deck_spec import normalize_deck_spec
from slide_smith.pptx_edit import delete_slide
from slide_smith.renderer import (
    FALLBACK_LAYOUT_ID,
    _layout_by_name,
    _make_fallback_slide_spec,
    _record_render_warning,
    _render_image_left_text_right,
    _render_section,
    _render_title,
    _render_title_and_bullets,
)
from slide_smith.styling import load_styles
from slide_smith.template_loader import load_template_spec


class EditError(Exception):
    """Raised when an edit operation cannot be completed safely."""



def _load_json(path: str) -> dict[str, Any]:
    return json.loads(Path(path).read_text())



def add_slide_to_deck(deck_path: str, after_index: int, archetype: str, input_path: str, template_id: str = "default") -> str:
    prs = Presentation(deck_path)
    slide_w_emu = int(prs.slide_width)
    slide_h_emu = int(prs.slide_height)
    if after_index != len(prs.slides) - 1:
        raise EditError(
            "Only append-style insert-slide is supported right now; use --after with the current last slide index"
        )

    template_spec = load_template_spec(template_id)
    styles = load_styles(template_spec)
    archetypes = {item["id"]: item for item in template_spec.get("archetypes", [])}

    raw_slide_spec = _load_json(input_path)
    if not isinstance(raw_slide_spec, dict):
        raise EditError("Slide input must be a JSON object")

    if archetype:
        raw_slide_spec.setdefault("layout_id", archetype)
    slide_spec, _warnings = normalize_deck_spec({"slides": [raw_slide_spec]})
    slide_spec = dict((slide_spec.get("slides") or [raw_slide_spec])[0])

    requested_layout_id = str(slide_spec.get("layout_id") or slide_spec.get("archetype") or archetype)
    effective_archetype = str(slide_spec.get("archetype") or requested_layout_id)
    base_dir = Path(input_path).resolve().parent

    def render_slide_for(kind: str, spec: dict[str, Any]) -> None:
        if kind not in archetypes:
            raise EditError(f"Layout '{kind}' is not supported by template '{template_id}'")

        layout_name = archetypes[kind]["layout"]
        slide = prs.slides.add_slide(_layout_by_name(prs, layout_name))
        archetype_spec = archetypes[kind]

        if kind == "title":
            _render_title(slide, spec, styles, archetype_spec, kind, slide_w_emu=slide_w_emu, slide_h_emu=slide_h_emu)
        elif kind == "section":
            _render_section(slide, spec, styles, archetype_spec, kind, slide_w_emu=slide_w_emu, slide_h_emu=slide_h_emu)
        elif kind in {"title_and_bullets", "title_subtitle_and_bullets"}:
            _render_title_and_bullets(slide, spec, styles, archetype_spec, kind, slide_w_emu=slide_w_emu, slide_h_emu=slide_h_emu)
        elif kind in {"image_left_text_right", "text_with_image"}:
            _render_image_left_text_right(slide, spec, base_dir, styles, archetype_spec, kind, slide_w_emu=slide_w_emu, slide_h_emu=slide_h_emu)
        else:
            raise EditError(f"Layout '{kind}' is not implemented")

    before_count = len(prs.slides)
    try:
        render_slide_for(effective_archetype, slide_spec)
    except EditError as exc:
        try:
            while len(prs.slides) > before_count:
                delete_slide(prs, len(prs.slides) - 1)
        except Exception:
            pass

        if requested_layout_id == FALLBACK_LAYOUT_ID:
            raise

        _record_render_warning(
            raw_slide_spec,
            slide_index=after_index + 1,
            requested_layout_id=requested_layout_id,
            reason=str(exc),
        )
        fallback_spec = _make_fallback_slide_spec(slide_spec, requested_layout_id=requested_layout_id, reason=str(exc))
        render_slide_for(FALLBACK_LAYOUT_ID, fallback_spec)

    prs.save(deck_path)
    return deck_path



def _resolve_image_path(image_field: object, base_dir: Path) -> Path:
    if isinstance(image_field, str):
        path = image_field
    elif isinstance(image_field, dict):
        path = image_field.get("path")
    else:
        raise EditError(f"Unsupported image field type: {type(image_field).__name__}")

    if not isinstance(path, str) or not path:
        raise EditError("Image replacement requires a non-empty path")

    p = Path(path).expanduser()
    resolved = (base_dir / p).resolve() if not p.is_absolute() else p.resolve()
    if not resolved.exists():
        raise EditError(f"Image file not found: {resolved}")
    if not resolved.is_file():
        raise EditError(f"Image path is not a file: {resolved}")
    return resolved



def update_slide_in_deck(deck_path: str, index: int, input_path: str) -> str:
    prs = Presentation(deck_path)
    if index < 0 or index >= len(prs.slides):
        raise EditError(f"Slide index {index} out of range")

    patch = _load_json(input_path)
    slide = prs.slides[index]

    if "title" in patch:
        try:
            slide.shapes.title.text = patch["title"]
        except Exception as exc:
            raise EditError("Target slide does not have a writable title placeholder") from exc

    if "subtitle" in patch:
        try:
            slide.placeholders[1].text = patch["subtitle"]
        except Exception as exc:
            raise EditError("Target slide does not support subtitle update via placeholder 1") from exc

    if "body" in patch:
        try:
            slide.placeholders[1].text = patch["body"]
        except Exception:
            try:
                slide.placeholders[2].text = patch["body"]
            except Exception as exc:
                raise EditError("Target slide does not support body update via known placeholders") from exc

    if "bullets" in patch:
        try:
            text_frame = slide.placeholders[1].text_frame
        except Exception as exc:
            raise EditError("Target slide does not support bullet updates") from exc
        bullets = patch["bullets"]
        if not isinstance(bullets, list):
            raise EditError("'bullets' patch must be an array of strings")
        if not bullets:
            text_frame.text = ""
        else:
            text_frame.text = str(bullets[0])
            for bullet in bullets[1:]:
                p = text_frame.add_paragraph()
                p.text = str(bullet)
                p.level = 0

    if "notes" in patch:
        notes = patch["notes"]
        if notes is not None and not isinstance(notes, str):
            raise EditError("'notes' patch must be a string (or null to clear)")
        slide.notes_slide.notes_text_frame.text = notes or ""

    if "image" in patch:
        base_dir = Path(input_path).resolve().parent
        resolved = _resolve_image_path(patch["image"], base_dir)
        try:
            image_placeholder = slide.placeholders[1]
        except Exception as exc:
            raise EditError("Target slide does not have image placeholder idx=1") from exc
        slide.shapes.add_picture(
            str(resolved),
            image_placeholder.left,
            image_placeholder.top,
            width=image_placeholder.width,
            height=image_placeholder.height,
        )

    prs.save(deck_path)
    return deck_path



def list_slides_in_deck(deck_path: str) -> list[dict[str, Any]]:
    prs = Presentation(deck_path)
    items: list[dict[str, Any]] = []
    for i, slide in enumerate(prs.slides):
        title = None
        try:
            title = slide.shapes.title.text
        except Exception:
            title = None
        items.append({"index": i, "title": title})
    return items



def delete_slide_from_deck(deck_path: str, index: int) -> str:
    prs = Presentation(deck_path)
    if index < 0 or index >= len(prs.slides):
        raise EditError(f"Slide index {index} out of range")
    delete_slide(prs, index)
    prs.save(deck_path)
    return deck_path
