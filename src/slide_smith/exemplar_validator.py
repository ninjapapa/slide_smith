from __future__ import annotations

import json
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation


def _enum_name(value: Any) -> str:
    name = getattr(value, "name", None)
    if isinstance(name, str) and name:
        return name
    return str(value)


def _layout_index_by_id(style_profile: dict[str, Any]) -> dict[str, int]:
    out: dict[str, int] = {}
    for l in style_profile.get("layouts") or []:
        if not isinstance(l, dict):
            continue
        lid = l.get("layoutId")
        idx = l.get("index")
        if isinstance(lid, str) and lid and isinstance(idx, int):
            out[lid] = idx
    return out


def _bbox(shape: Any) -> tuple[int, int, int, int]:
    left = int(getattr(shape, "left", 0) or 0)
    top = int(getattr(shape, "top", 0) or 0)
    width = int(getattr(shape, "width", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)
    return left, top, width, height


def _bbox_close(a: tuple[int, int, int, int], b: tuple[int, int, int, int], *, tol: int = 20) -> bool:
    return all(abs(x - y) <= tol for x, y in zip(a, b))


@dataclass(frozen=True)
class ValidateExemplarResult:
    ok: bool
    errors: list[str]
    warnings: list[str]
    report: dict[str, Any]


def validate_exemplar(
    *,
    reference_pptx: str,
    pptx: str,
    style_profile: dict[str, Any],
) -> ValidateExemplarResult:
    """Validate an output PPTX against a reference deck StyleProfile.

    v1.2 checks (deterministic, best-effort):
    - slide layouts are among the reference deck layouts (by index/name/schema)
    - slide edits are "placeholder-only" in spirit:
      - allow placeholder shapes
      - allow pictures occupying known picture placeholder bboxes

    Theme checks are not enforced in v1.2 (warn only; style extraction may be incomplete).
    """

    ref_path = Path(reference_pptx).expanduser().resolve()
    out_path = Path(pptx).expanduser().resolve()

    if not ref_path.exists() or not ref_path.is_file():
        raise FileNotFoundError(f"Reference PPTX not found: {ref_path}")
    if not out_path.exists() or not out_path.is_file():
        raise FileNotFoundError(f"PPTX not found: {out_path}")

    errors: list[str] = []
    warnings: list[str] = []

    ref = Presentation(str(ref_path))
    prs = Presentation(str(out_path))

    # Map of allowed layout indices from style_profile
    lid_to_idx = _layout_index_by_id(style_profile)
    allowed_layout_indices = set(lid_to_idx.values())

    # Build allowed picture placeholder bboxes per layout index (from reference)
    allowed_pic_bboxes: dict[int, list[tuple[int, int, int, int]]] = {}
    for l in style_profile.get("layouts") or []:
        if not isinstance(l, dict):
            continue
        idx = l.get("index")
        if not isinstance(idx, int):
            continue
        bxs: list[tuple[int, int, int, int]] = []
        for ph in l.get("placeholders") or []:
            if not isinstance(ph, dict):
                continue
            if str(ph.get("type", "")).lower() != "picture":
                continue
            bb = ph.get("bbox") or {}
            if not isinstance(bb, dict):
                continue
            try:
                bxs.append((int(bb.get("x")), int(bb.get("y")), int(bb.get("w")), int(bb.get("h"))))
            except Exception:
                continue
        if bxs:
            allowed_pic_bboxes[idx] = bxs

    # Validate slide-by-slide
    for i, slide in enumerate(prs.slides, start=1):
        layout = slide.slide_layout
        layout_name = getattr(layout, "name", "")

        # Determine layout index within *reference* presentation by matching name and placeholder schema.
        # Best-effort: first attempt to find same object index by name.
        matched_ref_idx: int | None = None

        # Try exact name match + placeholder signature match
        out_sig = [(int(ph.placeholder_format.idx), _enum_name(ph.placeholder_format.type)) for ph in layout.placeholders]

        for ridx, rlayout in enumerate(ref.slide_layouts):
            if rlayout.name != layout_name:
                continue
            ref_sig = [(int(ph.placeholder_format.idx), _enum_name(ph.placeholder_format.type)) for ph in rlayout.placeholders]
            if ref_sig == out_sig:
                matched_ref_idx = ridx
                break

        if matched_ref_idx is None:
            # As a fallback, accept if layout name exists in reference at all.
            name_exists = any(rlayout.name == layout_name for rlayout in ref.slide_layouts)
            if not name_exists:
                errors.append(f"slide[{i}]: layout '{layout_name}' not found in reference")
            else:
                warnings.append(f"slide[{i}]: layout '{layout_name}' signature mismatch; validation is best-effort")
        else:
            if allowed_layout_indices and matched_ref_idx not in allowed_layout_indices:
                errors.append(
                    f"slide[{i}]: layout '{layout_name}' index {matched_ref_idx} not present in style_profile.layouts"
                )

        # Placeholder-only edits check (best-effort)
        pic_boxes = allowed_pic_bboxes.get(matched_ref_idx or -1, [])
        for sh in slide.shapes:
            # If it's a placeholder, accept.
            if getattr(sh, "is_placeholder", False):
                continue
            try:
                _ = sh.placeholder_format  # may raise if not a placeholder
                continue
            except Exception:
                pass

            st = _enum_name(getattr(sh, "shape_type", None))
            if st == "PICTURE" and pic_boxes:
                bb = _bbox(sh)
                if any(_bbox_close(bb, pb) for pb in pic_boxes):
                    continue

            # Otherwise, unexpected shape.
            errors.append(
                f"slide[{i}]: unexpected shape '{getattr(sh, 'name', '')}' type={st} bbox={_bbox(sh)}"
            )

    report = {
        "ok": len(errors) == 0,
        "reference": str(ref_path),
        "pptx": str(out_path),
        "errors": errors,
        "warnings": warnings,
    }
    json.dumps(report, sort_keys=True)
    return ValidateExemplarResult(ok=(len(errors) == 0), errors=errors, warnings=warnings, report=report)
