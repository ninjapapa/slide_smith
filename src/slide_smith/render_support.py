from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation

from slide_smith.styling import apply_text_style


class RenderingError(Exception):
    """Raised when a deck cannot be rendered."""


def _layout_by_name(prs: Presentation, name: str):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    raise RenderingError(f"Slide layout '{name}' not found in template presentation")


def _layout_for_archetype(prs: Presentation, archetype_spec: dict[str, Any]):
    """Resolve a layout for an archetype.

    Prefer `layout_part` when available (more stable for rich branded templates).
    Fallback to `layout` name.
    """

    from slide_smith.layout_resolver import LayoutResolveError, resolve_layout

    layout_name = archetype_spec.get("layout")
    layout_part = archetype_spec.get("layout_part")

    try:
        return resolve_layout(prs=prs, layout_name=layout_name, layout_part=layout_part).layout
    except LayoutResolveError as exc:
        raise RenderingError(str(exc)) from exc


def _slot_spec(archetype_spec: dict[str, Any], slot_name: str) -> dict[str, Any] | None:
    for slot in archetype_spec.get("slots", []):
        if isinstance(slot, dict) and slot.get("name") == slot_name:
            return slot
    return None


def _slot_index(archetype_spec: dict[str, Any], slot_name: str, default: int | None = None) -> int | None:
    """Resolve placeholder_idx for a slot.

    `default` is only used when the slot is missing entirely.
    If the slot exists but does not declare placeholder_idx (e.g. box-based slots),
    this returns None.
    """

    slot = _slot_spec(archetype_spec, slot_name)
    if slot is None:
        return default
    idx = slot.get("placeholder_idx")
    if isinstance(idx, int):
        return idx
    return None


def _slot_box(archetype_spec: dict[str, Any], slot_name: str) -> dict[str, Any] | None:
    slot = _slot_spec(archetype_spec, slot_name)
    if slot is None:
        return None
    box = slot.get("box")
    if isinstance(box, dict):
        return box
    return None


def _box_to_emu(box: dict[str, Any], *, slide_w_emu: int, slide_h_emu: int) -> tuple[int, int, int, int] | None:
    """Convert a slot box to EMU coordinates.

    Supported forms:
    - units=relative: x,y,w,h are floats in [0,1]
    - units=emu: left,top,width,height are ints (or x,y,w,h)

    Returns: (left, top, width, height)
    """

    units = box.get("units")
    if units == "relative":
        try:
            x = float(box["x"])
            y = float(box["y"])
            w = float(box["w"])
            h = float(box["h"])
        except Exception:
            return None
        left = int(x * slide_w_emu)
        top = int(y * slide_h_emu)
        width = int(w * slide_w_emu)
        height = int(h * slide_h_emu)
        return left, top, width, height

    if units == "emu":
        def gi(*keys):
            for k in keys:
                v = box.get(k)
                if isinstance(v, int):
                    return v
            return None

        left = gi("left", "x")
        top = gi("top", "y")
        width = gi("width", "w")
        height = gi("height", "h")
        if None in (left, top, width, height):
            return None
        return int(left), int(top), int(width), int(height)

    return None


def _required_slot_target(
    archetype_id: str,
    archetype_spec: dict[str, Any],
    slot_name: str,
    *,
    default_idx: int | None = None,
) -> tuple[int | None, dict[str, Any] | None]:
    slot = _slot_spec(archetype_spec, slot_name)
    idx = _slot_index(archetype_spec, slot_name, default_idx)
    box = _slot_box(archetype_spec, slot_name)

    required = bool(slot.get("required")) if slot is not None else False
    if required and idx is None and box is None:
        raise RenderingError(
            f"Template archetype '{archetype_id}' missing required slot mapping '{slot_name}' (placeholder_idx or box)"
        )

    return idx, box


def _set_placeholder_text(slide, idx: int | None, text: str | None, style=None, *, context: str = "") -> None:
    if idx is None or text is None:
        return
    try:
        ph = slide.placeholders[idx]
        ph.text = text
        if style is not None and hasattr(ph, "text_frame"):
            apply_text_style(ph.text_frame, style)
    except KeyError as exc:
        ctx = f" ({context})" if context else ""
        raise RenderingError(f"Placeholder idx={idx} not found on slide{ctx}") from exc


def _set_box_text(slide, box_emu: tuple[int, int, int, int] | None, text: str | None, style=None) -> None:
    if box_emu is None or text is None:
        return
    left, top, width, height = box_emu
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.text = text
    if style is not None:
        apply_text_style(tf, style)


def _set_box_bullets(slide, box_emu: tuple[int, int, int, int] | None, bullets: list[str], style=None) -> None:
    if box_emu is None:
        return
    left, top, width, height = box_emu
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    if not bullets:
        tf.text = ""
        return
    tf.text = bullets[0]
    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
    if style is not None:
        apply_text_style(tf, style)


def _set_box_image(slide, box_emu: tuple[int, int, int, int] | None, image_path: Path | None) -> None:
    if box_emu is None or image_path is None:
        return
    left, top, width, height = box_emu
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def _set_slot_text(
    slide,
    archetype_id: str,
    archetype_spec: dict[str, Any],
    slot_name: str,
    text: str | None,
    style,
    *,
    slide_w_emu: int,
    slide_h_emu: int,
    default_idx: int | None = None,
    context: str = "",
) -> None:
    idx, box = _required_slot_target(archetype_id, archetype_spec, slot_name, default_idx=default_idx)
    if idx is not None:
        _set_placeholder_text(slide, idx, text, style, context=context)
        return
    if box is not None:
        box_emu = _box_to_emu(box, slide_w_emu=slide_w_emu, slide_h_emu=slide_h_emu)
        _set_box_text(slide, box_emu, text, style)


def _resolve_image_path(base_dir: Path, image_field: Any) -> Path | None:
    image_path = None
    if isinstance(image_field, str):
        image_path = image_field
    elif isinstance(image_field, dict):
        image_path = image_field.get("path")

    if not image_path:
        return None

    p = Path(str(image_path)).expanduser()
    resolved = (base_dir / p).resolve() if not p.is_absolute() else p.resolve()
    if not resolved.exists():
        raise RenderingError(f"Image file not found: {resolved}")
    if not resolved.is_file():
        raise RenderingError(f"Image path is not a file: {resolved}")
    return resolved


def _set_slot_image(
    slide,
    archetype_id: str,
    archetype_spec: dict[str, Any],
    slot_name: str,
    image_path: Path | None,
    *,
    slide_w_emu: int,
    slide_h_emu: int,
    default_idx: int | None = None,
    context: str = "",
) -> None:
    if image_path is None:
        return

    idx, box = _required_slot_target(archetype_id, archetype_spec, slot_name, default_idx=default_idx)
    if idx is not None:
        try:
            image_placeholder = slide.placeholders[idx]
        except KeyError as exc:
            ctx = f" ({context})" if context else ""
            raise RenderingError(f"Placeholder idx={idx} not found on slide{ctx}") from exc
        slide.shapes.add_picture(
            str(image_path),
            image_placeholder.left,
            image_placeholder.top,
            width=image_placeholder.width,
            height=image_placeholder.height,
        )
        return

    if box is not None:
        box_emu = _box_to_emu(box, slide_w_emu=slide_w_emu, slide_h_emu=slide_h_emu)
        _set_box_image(slide, box_emu, image_path)
