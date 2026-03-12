from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation


@dataclass(frozen=True)
class InspectPptxResult:
    pptx: str
    slide_size: dict[str, int]
    layouts: list[dict[str, Any]]


@dataclass(frozen=True)
class InspectSlideResult:
    pptx: str
    slide_number: int
    slide_size: dict[str, int]
    shapes: list[dict[str, Any]]


def _enum_name(value: Any) -> str:
    """Best-effort conversion of python-pptx enum-ish values into stable strings."""

    # Many python-pptx enums are `EnumValue` with `.name`.
    name = getattr(value, "name", None)
    if isinstance(name, str) and name:
        return name
    return str(value)


def inspect_pptx(pptx_path: str) -> InspectPptxResult:
    path = Path(pptx_path).expanduser()
    if not path.exists():
        raise FileNotFoundError(f"PPTX not found: {path}")
    if not path.is_file():
        raise FileNotFoundError(f"PPTX path is not a file: {path}")

    abs_path = str(path.resolve())
    prs = Presentation(abs_path)

    slide_size = {
        "width_emu": int(prs.slide_width),
        "height_emu": int(prs.slide_height),
    }

    layouts: list[dict[str, Any]] = []
    for idx, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for ph in sorted(layout.placeholders, key=lambda p: int(p.placeholder_format.idx)):
            placeholders.append(
                {
                    "idx": int(ph.placeholder_format.idx),
                    "name": getattr(ph, "name", ""),
                    "ph_type": _enum_name(ph.placeholder_format.type),
                    "shape_type": _enum_name(getattr(ph, "shape_type", None)),
                }
            )

        layouts.append(
            {
                "name": layout.name,
                "index": idx,
                "placeholders": placeholders,
            }
        )

    return InspectPptxResult(pptx=abs_path, slide_size=slide_size, layouts=layouts)


def inspect_slide(pptx_path: str, *, slide_number: int) -> InspectSlideResult:
    """Inspect a single slide instance (1-indexed) and return shape inventory.

    This is intentionally "agent-friendly": include geometry (EMU and normalized),
    shape type, placeholder metadata when present, and a short text preview.
    """

    path = Path(pptx_path).expanduser()
    if not path.exists():
        raise FileNotFoundError(f"PPTX not found: {path}")
    if not path.is_file():
        raise FileNotFoundError(f"PPTX path is not a file: {path}")

    abs_path = str(path.resolve())
    prs = Presentation(abs_path)

    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError(f"slide_number out of range: {slide_number} (1..{len(prs.slides)})")

    slide = prs.slides[slide_number - 1]

    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)

    def rel(v: int, denom: int) -> float:
        if denom <= 0:
            return 0.0
        return float(v) / float(denom)

    shapes: list[dict[str, Any]] = []
    for z, shape in enumerate(slide.shapes):
        left = int(getattr(shape, "left", 0) or 0)
        top = int(getattr(shape, "top", 0) or 0)
        width = int(getattr(shape, "width", 0) or 0)
        height = int(getattr(shape, "height", 0) or 0)

        item: dict[str, Any] = {
            "z": int(z),
            "name": getattr(shape, "name", ""),
            "shape_type": _enum_name(getattr(shape, "shape_type", None)),
            "bbox_emu": {"left": left, "top": top, "width": width, "height": height},
            "bbox_rel": {
                "x": rel(left, slide_w),
                "y": rel(top, slide_h),
                "w": rel(width, slide_w),
                "h": rel(height, slide_h),
            },
        }

        # Placeholder metadata (if any)
        phf = getattr(shape, "placeholder_format", None)
        if phf is not None:
            try:
                item["placeholder"] = {
                    "idx": int(getattr(phf, "idx")),
                    "ph_type": _enum_name(getattr(phf, "type", None)),
                    "name": getattr(shape, "name", ""),
                }
            except Exception:
                pass

        # Text preview
        text = None
        if getattr(shape, "has_text_frame", False):
            try:
                raw = shape.text_frame.text or ""
                raw = " ".join(raw.split())
                text = raw[:200] if raw else None
            except Exception:
                text = None
        if text:
            item["text"] = text

        # Picture metadata (best-effort)
        img = getattr(shape, "image", None)
        if img is not None:
            try:
                item["image"] = {
                    "content_type": getattr(img, "content_type", None),
                    "filename": getattr(img, "filename", None),
                    "size": int(getattr(img, "size")),
                }
            except Exception:
                pass

        shapes.append(item)

    return InspectSlideResult(
        pptx=abs_path,
        slide_number=int(slide_number),
        slide_size={"width_emu": slide_w, "height_emu": slide_h},
        shapes=shapes,
    )
