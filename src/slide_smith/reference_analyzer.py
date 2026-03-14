from __future__ import annotations

import hashlib
import json
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation


def _enum_name(value: Any) -> str:
    name = getattr(value, "name", None)
    if isinstance(name, str) and name:
        return name
    return str(value)


def _sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def _slug(s: str) -> str:
    s = s.strip().lower()
    s = re.sub(r"[^a-z0-9]+", "-", s)
    s = re.sub(r"-+", "-", s).strip("-")
    return s or "layout"


def _stable_layout_id(
    *,
    layout_index: int,
    layout_name: str,
    placeholders: list[dict[str, Any]],
    part: str | None = None,
) -> str:
    """Generate a stable-ish ID for a layout for a given reference deck.

    Prefer identifying layouts by OpenXML part name when available (raw mode),
    because indices and names can shift.

    Derived from:
    - (optional) slideLayout part name, e.g. `ppt/slideLayouts/slideLayout12.xml`
    - layout name
    - placeholder signature (type+idx+bbox)

    This should remain stable for the same deck unless layouts change.
    """

    signature = {
        "part": str(part or ""),
        "name": str(layout_name),
        "placeholders": [
            {
                "type": str(p.get("type", "")),
                "idx": int(p.get("idx", -1)),
                "bbox": p.get("bbox"),
            }
            for p in placeholders
        ],
    }
    raw = json.dumps(signature, sort_keys=True, separators=(",", ":")).encode("utf-8")
    digest = hashlib.sha1(raw).hexdigest()  # stable + short; not for security

    if part:
        part_slug = _slug(Path(part).stem)
        return f"layout:{part_slug}:{digest[:10]}"

    return f"layout:{layout_index}:{_slug(layout_name)}:{digest[:10]}"


@dataclass(frozen=True)
class AnalyzeReferenceResult:
    style_profile: dict[str, Any]


def analyze_reference(pptx_path: str, *, mode: str = "pptx") -> AnalyzeReferenceResult:
    path = Path(pptx_path).expanduser()
    if not path.exists():
        raise FileNotFoundError(f"PPTX not found: {path}")
    if not path.is_file():
        raise FileNotFoundError(f"PPTX path is not a file: {path}")

    abs_path = path.resolve()
    sha256 = _sha256_file(abs_path)

    mode = (mode or "pptx").strip().lower()
    if mode not in ("pptx", "raw"):
        raise ValueError(f"Unsupported mode: {mode} (expected 'pptx' or 'raw')")

    # Slide size: use python-pptx in pptx mode; raw OpenXML in raw mode.
    if mode == "pptx":
        prs = Presentation(str(abs_path))
        slide_size = {"widthEmu": int(prs.slide_width), "heightEmu": int(prs.slide_height)}
    else:
        from slide_smith.openxml_presentation import inspect_openxml_presentation

        pres = inspect_openxml_presentation(str(abs_path))
        slide_size = {"widthEmu": int(pres.slide_size["width_emu"]), "heightEmu": int(pres.slide_size["height_emu"])}

    layouts: list[dict[str, Any]] = []

    if mode == "pptx":
        # prs defined above in pptx mode
        for idx, layout in enumerate(prs.slide_layouts):
            placeholders: list[dict[str, Any]] = []

            # Note: layout.placeholders are placeholder shapes on the layout.
            # They have geometry (left/top/width/height) in EMU.
            for ph in sorted(layout.placeholders, key=lambda p: int(p.placeholder_format.idx)):
                left = int(getattr(ph, "left", 0) or 0)
                top = int(getattr(ph, "top", 0) or 0)
                width = int(getattr(ph, "width", 0) or 0)
                height = int(getattr(ph, "height", 0) or 0)

                placeholders.append(
                    {
                        "type": _enum_name(ph.placeholder_format.type),
                        "idx": int(ph.placeholder_format.idx),
                        "name": getattr(ph, "name", ""),
                        "shapeType": _enum_name(getattr(ph, "shape_type", None)),
                        "bbox": {"x": left, "y": top, "w": width, "h": height},
                    }
                )

            layout_name = getattr(layout, "name", "") or f"Layout {idx}"
            layout_id = _stable_layout_id(layout_index=idx, layout_name=layout_name, placeholders=placeholders)

            layouts.append(
                {
                    "layoutId": layout_id,
                    "name": layout_name,
                    "index": int(idx),
                    "placeholders": placeholders,
                }
            )

    else:
        # raw openxml mode: enumerate ppt/slideLayouts/ parts.
        from slide_smith.openxml_layouts import inspect_openxml_layouts

        raw = inspect_openxml_layouts(str(abs_path))
        for idx, layout in enumerate(raw.layouts):
            part = str(layout.get("part", ""))
            placeholders = []
            for ph in layout.get("placeholders") or []:
                # raw types are already strings
                placeholders.append(
                    {
                        "type": str(ph.get("type", "")),
                        "idx": int(ph.get("idx", -1)),
                        "name": "",
                        "shapeType": "",
                        "bbox": ph.get("bbox") or {"x": 0, "y": 0, "w": 0, "h": 0},
                    }
                )

            layout_name = str(layout.get("name") or f"Layout {idx}")
            layout_id = _stable_layout_id(
                layout_index=idx,
                layout_name=layout_name,
                placeholders=placeholders,
                part=part or None,
            )
            layouts.append(
                {
                    "layoutId": layout_id,
                    "name": layout_name,
                    "index": int(idx),
                    "part": part,
                    "placeholders": placeholders,
                }
            )

    # Theme extraction: best-effort. python-pptx doesn't expose full theme reliably.
    # Keep this as a placeholder structure for now.
    theme: dict[str, Any] = {}

    style_profile: dict[str, Any] = {
        "version": "1",
        "reference": {
            "path": str(abs_path),
            "sha256": sha256,
            "slideSize": slide_size,
        },
        "theme": theme,
        "layouts": layouts,
        "constraints": {"placeholderOnly": True, "allowNewShapes": False},
    }

    return AnalyzeReferenceResult(style_profile=style_profile)
