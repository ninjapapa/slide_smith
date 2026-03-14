from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation


class BootstrapFromSlideError(Exception):
    """Raised when slide-instance bootstrap cannot be completed."""


@dataclass(frozen=True)
class BootstrappedArchetype:
    pptx: str
    slide_number: int
    slide_layout: str
    archetype_id: str
    archetype_spec: dict[str, Any]


def _norm_box(*, left: int, top: int, width: int, height: int, slide_w: int, slide_h: int) -> dict[str, Any]:
    if slide_w <= 0 or slide_h <= 0:
        raise BootstrapFromSlideError("Invalid slide dimensions")
    return {
        "units": "relative",
        "x": left / slide_w,
        "y": top / slide_h,
        "w": width / slide_w,
        "h": height / slide_h,
    }


def bootstrap_archetype_from_slide(
    pptx_path: str,
    *,
    slide_number: int,
    archetype_id: str,
    prefer_placeholders: bool = True,
) -> BootstrappedArchetype:
    """Infer a box-based archetype spec from a concrete slide instance.

    MVP heuristics:
    - title: top-most non-empty text box
    - body: next top-most non-empty text box
    - image: largest picture by area

    Returns an archetype spec that uses `box` mappings (units=relative).
    """

    p = Path(pptx_path).expanduser()
    if not p.exists() or not p.is_file():
        raise BootstrapFromSlideError(f"PPTX not found: {p}")

    prs = Presentation(str(p))

    if slide_number < 1 or slide_number > len(prs.slides):
        raise BootstrapFromSlideError(f"slide_number out of range: {slide_number} (1..{len(prs.slides)})")

    slide = prs.slides[slide_number - 1]
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)

    text_shapes: list[dict[str, Any]] = []
    picture_shapes: list[dict[str, Any]] = []

    for shape in slide.shapes:
        left = int(getattr(shape, "left", 0) or 0)
        top = int(getattr(shape, "top", 0) or 0)
        width = int(getattr(shape, "width", 0) or 0)
        height = int(getattr(shape, "height", 0) or 0)
        area = width * height

        if getattr(shape, "has_text_frame", False):
            raw = (shape.text_frame.text or "").strip()  # type: ignore[attr-defined]
            if raw:
                text_shapes.append(
                    {
                        "left": left,
                        "top": top,
                        "width": width,
                        "height": height,
                        "area": area,
                        "text": " ".join(raw.split())[:200],
                    }
                )

        img = getattr(shape, "image", None)
        if img is not None:
            picture_shapes.append(
                {
                    "left": left,
                    "top": top,
                    "width": width,
                    "height": height,
                    "area": area,
                }
            )

    text_shapes.sort(key=lambda s: (s["top"], s["left"]))
    picture_shapes.sort(key=lambda s: s["area"], reverse=True)

    def pick_title_box() -> dict[str, Any] | None:
        if not text_shapes:
            return None
        s = text_shapes[0]
        return _norm_box(left=s["left"], top=s["top"], width=s["width"], height=s["height"], slide_w=slide_w, slide_h=slide_h)

    def pick_body_box() -> dict[str, Any] | None:
        if len(text_shapes) < 2:
            return None
        s = text_shapes[1]
        return _norm_box(left=s["left"], top=s["top"], width=s["width"], height=s["height"], slide_w=slide_w, slide_h=slide_h)

    def pick_image_box() -> dict[str, Any] | None:
        if not picture_shapes:
            return None
        s = picture_shapes[0]
        return _norm_box(left=s["left"], top=s["top"], width=s["width"], height=s["height"], slide_w=slide_w, slide_h=slide_h)

    title_box = pick_title_box()
    body_box = pick_body_box()
    image_box = pick_image_box()

    def placeholder_idx_by_type(*want_types: set[int]) -> int | None:
        # Best-effort placeholder discovery on the concrete slide.
        for ph in getattr(slide, "placeholders", []):
            try:
                pht = int(ph.placeholder_format.type)  # type: ignore[attr-defined]
                idx = int(ph.placeholder_format.idx)  # type: ignore[attr-defined]
            except Exception:
                continue
            if pht in want_types:
                return idx
        return None

    if archetype_id == "image_left_text_right":
        # Prefer real placeholders when possible (placeholder-first).
        title_idx = None
        body_idx = None
        image_idx = None

        if prefer_placeholders:
            try:
                from pptx.enum.shapes import PP_PLACEHOLDER  # type: ignore

                title_idx = placeholder_idx_by_type({int(PP_PLACEHOLDER.TITLE), int(PP_PLACEHOLDER.CENTER_TITLE)})
                body_idx = placeholder_idx_by_type({int(PP_PLACEHOLDER.BODY)})
                image_idx = placeholder_idx_by_type({int(PP_PLACEHOLDER.PICTURE)})
            except Exception:
                title_idx = body_idx = image_idx = None

        slots: list[dict[str, Any]] = []
        heuristic = []

        if title_idx is not None and body_idx is not None and image_idx is not None:
            heuristic.append("placeholder-first: TITLE/BODY/PICTURE placeholders")
            slots = [
                {"name": "title", "type": "text", "required": True, "placeholder_idx": int(title_idx)},
                {"name": "image", "type": "image", "required": True, "placeholder_idx": int(image_idx)},
                {"name": "body", "type": "text", "required": True, "placeholder_idx": int(body_idx)},
            ]
        else:
            heuristic.append("box-fallback: topmost text=title, next text=body, largest picture=image")
            if title_box is None:
                raise BootstrapFromSlideError("Could not infer title box (no non-empty text shapes)")
            if body_box is None:
                raise BootstrapFromSlideError("Could not infer body box (need at least 2 non-empty text shapes)")
            if image_box is None:
                raise BootstrapFromSlideError("Could not infer image box (no picture shapes)")
            slots = [
                {"name": "title", "type": "text", "required": True, "box": title_box},
                {"name": "image", "type": "image", "required": True, "box": image_box},
                {"name": "body", "type": "text", "required": True, "box": body_box},
            ]

        archetype_spec = {
            "id": archetype_id,
            "layout": slide.slide_layout.name,
            "description": f"Bootstrapped from slide {slide_number} of {p.name}",
            "bootstrap": {
                "source": {"pptx": str(p.resolve()), "slide_number": int(slide_number)},
                "heuristic": "; ".join(heuristic),
                "prefer_placeholders": bool(prefer_placeholders),
            },
            "slots": slots,
        }
        return BootstrappedArchetype(
            pptx=str(p.resolve()),
            slide_number=int(slide_number),
            slide_layout=slide.slide_layout.name,
            archetype_id=archetype_id,
            archetype_spec=archetype_spec,
        )

    raise BootstrapFromSlideError(f"archetype '{archetype_id}' not supported by bootstrap-from-slide yet")
