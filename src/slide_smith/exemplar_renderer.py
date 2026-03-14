from __future__ import annotations

import json
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation


@dataclass(frozen=True)
class RenderExemplarResult:
    output_pptx: str


def _remove_all_slides(prs: Presentation) -> None:
    """Remove all slides from a presentation (python-pptx has no public API).

    We must remove both:
    - the slide id entry from the slide id list
    - the relationship from the presentation part

    Otherwise python-pptx may emit duplicate slide part names on save.
    """

    sldIdLst = prs.slides._sldIdLst  # type: ignore[attr-defined]
    # Iterate a copy; remove from XML list and drop rels.
    for sldId in list(sldIdLst):
        rId = sldId.rId
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def _layout_index_by_id(style_profile: dict[str, Any]) -> dict[str, int]:
    out: dict[str, int] = {}
    for l in style_profile.get("layouts") or []:
        if not isinstance(l, dict):
            continue
        lid = l.get("layoutId")
        idx = l.get("index")
        if isinstance(lid, str) and lid and isinstance(idx, int):
            out[lid] = idx
    return out


def _fill_text_placeholder(slide: Any, *, ph_type: str, ph_idx: int, text: str) -> bool:
    """Fill a placeholder on a slide by type+idx. Returns True if filled."""

    for shape in slide.placeholders:
        try:
            pf = shape.placeholder_format
            if int(pf.idx) != int(ph_idx):
                continue
            if str(pf.type.name) != str(ph_type) and str(pf.type) != str(ph_type):
                # profile uses enum name; try both
                continue
        except Exception:
            continue

        if not getattr(shape, "has_text_frame", False):
            continue

        shape.text_frame.clear()
        # Keep simple: set entire text. (Bullets can be expanded later.)
        shape.text_frame.text = text
        return True

    return False


def _fill_picture_placeholder(slide: Any, *, ph_type: str, ph_idx: int, image_path: Path) -> bool:
    for shape in slide.placeholders:
        try:
            pf = shape.placeholder_format
            if int(pf.idx) != int(ph_idx):
                continue
            if str(pf.type.name) != str(ph_type) and str(pf.type) != str(ph_type):
                continue
        except Exception:
            continue

        # python-pptx picture placeholders provide insert_picture
        insert = getattr(shape, "insert_picture", None)
        if callable(insert):
            insert(str(image_path))
            return True

    return False


def render_exemplar(
    *,
    reference_pptx: str,
    style_profile: dict[str, Any],
    deck_spec: dict[str, Any],
    output_pptx: str,
    assets_base_dir: str | None = None,
) -> RenderExemplarResult:
    """Render exemplar-first deck_spec by instantiating layouts from reference PPTX.

    Deterministic constraints (v1.2):
    - Only create slides using layouts present in style_profile.
    - Only fill placeholders (no arbitrary new shapes).
    """

    ref = Path(reference_pptx).expanduser()
    if not ref.exists() or not ref.is_file():
        raise FileNotFoundError(f"Reference PPTX not found: {ref}")

    out_path = Path(output_pptx).expanduser().resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)

    lid_to_idx = _layout_index_by_id(style_profile)

    prs = Presentation(str(ref.resolve()))
    _remove_all_slides(prs)

    slides = deck_spec.get("slides") or []
    if not isinstance(slides, list) or not slides:
        raise ValueError("deck_spec.slides must be a non-empty list")

    base = Path(assets_base_dir).expanduser().resolve() if assets_base_dir else Path.cwd()

    for i, s in enumerate(slides):
        if not isinstance(s, dict):
            raise TypeError(f"deck_spec.slides[{i}] must be dict")

        layout_id = s.get("layoutId")
        if not isinstance(layout_id, str) or not layout_id:
            raise ValueError(f"deck_spec.slides[{i}].layoutId must be string")

        if layout_id not in lid_to_idx:
            raise ValueError(f"Unknown layoutId in deck_spec: {layout_id}")

        layout_idx = lid_to_idx[layout_id]
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)

        for f in s.get("fills") or []:
            if not isinstance(f, dict):
                continue
            ph = f.get("placeholder") or {}
            if not isinstance(ph, dict):
                continue

            ph_type = ph.get("type")
            ph_idx = ph.get("idx")
            if not isinstance(ph_type, str) or not isinstance(ph_idx, int):
                continue

            if isinstance(f.get("text"), str):
                ok = _fill_text_placeholder(slide, ph_type=ph_type, ph_idx=ph_idx, text=str(f["text"]))
                if not ok:
                    raise ValueError(f"Failed to fill text placeholder type={ph_type} idx={ph_idx} on slide {i}")

            if isinstance(f.get("image"), dict) and isinstance(f["image"].get("path"), str):
                img_path = (base / f["image"]["path"]).resolve() if not Path(f["image"]["path"]).is_absolute() else Path(f["image"]["path"]).expanduser().resolve()
                if not img_path.exists():
                    raise FileNotFoundError(f"Image not found: {img_path}")
                ok = _fill_picture_placeholder(slide, ph_type=ph_type, ph_idx=ph_idx, image_path=img_path)
                if not ok:
                    raise ValueError(f"Failed to fill picture placeholder type={ph_type} idx={ph_idx} on slide {i}")

    prs.save(str(out_path))

    # Ensure it is writable and JSON-serializable for callers
    json.dumps({"output_pptx": str(out_path)})
    return RenderExemplarResult(output_pptx=str(out_path))
