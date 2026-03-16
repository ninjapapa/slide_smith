from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation

from slide_smith.template_loader import load_template_spec, template_dir


@dataclass(frozen=True)
class TemplateValidationResult:
    ok: bool
    errors: list[str]


class TemplateValidationError(Exception):
    pass


def _validate_semantic(archetypes: list[object], profile: str) -> list[str]:
    errors: list[str] = []

    conventions_doc = "docs/layout-ids.md (Template slot conventions)"

    if profile not in {"standard", "extended"}:
        return errors

    required: dict[str, dict[str, bool]] = {
        # Base (standard profile)
        "title": {"title": True},
        "section": {"title": True},
        "title_and_bullets": {"title": True, "bullets": True},
        "title_subtitle_and_bullets": {"title": True, "subtitle": True, "bullets": True},
        "text_with_image": {"title": True, "image": True, "body": True},

    }

    if profile == "extended":
        required |= {
            "two_col": {"title": True, "col1_body": True, "col2_body": True},
            "version_page": {"title": True, "table_text": True},

            # agenda_with_image supports two template approaches:
            # - a single bullets/body box (slot=bullets)
            # - dedicated item placeholders (slot=item1_body and optional item1_marker)
            "agenda_with_image": {"title": True, "image": True, "bullets": False, "item1_body": False},

            "three_col_with_icons": {
                "title": True,
                "col1_title": True,
                "col1_body": True,
                "col1_icon": True,
                "col2_title": True,
                "col2_body": True,
                "col2_icon": True,
                "col3_title": True,
                "col3_body": True,
                "col3_icon": True,
            },
            "picture_compare": {
                "title": True,
                "left_image": True,
                "right_image": True,
                # compare text is optional
                "left_body": False,
                "right_body": False,
                "left_title": False,
                "right_title": False,
            },
        }

    typed_archetypes = [a for a in archetypes if isinstance(a, dict)]
    by_id = {a.get("id"): a for a in typed_archetypes}

    prefix = f"{profile} profile"
    for aid, req_slots in required.items():
        if aid not in by_id:
            errors.append(f"{prefix}: missing required layout id '{aid}'")
            continue
        a = by_id[aid]
        slots = a.get("slots") or []
        if not isinstance(slots, list):
            errors.append(f"{prefix}: layout id '{aid}' slots must be a list")
            continue
        slot_by_name = {s.get("name"): s for s in slots if isinstance(s, dict)}

        # Special-case: agenda_with_image can be satisfied by either a bullets slot
        # OR dedicated item placeholders.
        if aid == "agenda_with_image":
            has_bullets = "bullets" in slot_by_name
            has_items = "item1_body" in slot_by_name
            if not (has_bullets or has_items):
                errors.append(
                    f"{prefix}: layout id 'agenda_with_image' must define either slot 'bullets' (fallback) "
                    f"or slot 'item1_body' (preferred item placeholders). See {conventions_doc}"
                )

        for slot_name, must in req_slots.items():
            if not must:
                continue
            if slot_name not in slot_by_name:
                extra = ""
                if aid in {"agenda_with_image", "three_col_with_icons", "five_col_with_icons", "picture_compare"}:
                    extra = f" See {conventions_doc}"
                errors.append(f"{prefix}: layout id '{aid}' missing required slot '{slot_name}'.{extra}")
                continue
            slot = slot_by_name[slot_name]
            idx = slot.get("placeholder_idx")
            box = slot.get("box")
            if not isinstance(idx, int) and not isinstance(box, dict):
                errors.append(
                    f"{prefix}: layout id '{aid}' slot '{slot_name}' missing placeholder_idx (int) or box (object)"
                )

    return errors


def _validate_structural(archetypes: list[object], prs: Presentation, *, pptx_path: str) -> list[str]:
    errors: list[str] = []

    layout_names = {layout.name for layout in prs.slide_layouts}

    def layout_by_name(name: str):
        for layout in prs.slide_layouts:
            if layout.name == name:
                return layout
        return None

    # Raw OpenXML inventory (for part-based matching when names don't resolve)
    raw_by_part: dict[str, dict[str, Any]] = {}
    try:
        from slide_smith.openxml_layouts import inspect_openxml_layouts

        raw = inspect_openxml_layouts(pptx_path)
        for l in raw.layouts:
            part = l.get("part")
            if isinstance(part, str) and part:
                raw_by_part[part] = l
    except Exception:
        raw_by_part = {}

    for a in archetypes:
        if not isinstance(a, dict):
            errors.append("archetypes[] must be objects")
            continue
        aid = a.get("id")
        layout_name = a.get("layout")
        if not isinstance(aid, str) or not aid:
            errors.append("layout definition missing required 'id'")
            continue
        # Box-only archetypes (slide-instance-derived) may not require a resolvable layout.
        slots = a.get("slots") or []
        box_only = False
        if isinstance(slots, list) and slots:
            typed_slots = [s for s in slots if isinstance(s, dict)]
            if typed_slots and all(("placeholder_idx" not in s) and isinstance(s.get("box"), dict) for s in typed_slots):
                box_only = True

        if not isinstance(layout_name, str) or not layout_name:
            if box_only:
                errors.append(f"warning: layout id '{aid}' missing layout (box-only definition; skipping layout checks)")
                continue
            errors.append(f"layout id '{aid}' missing required 'layout'")
            continue

        layout_part = a.get("layout_part")
        if isinstance(layout_part, str) and layout_part.startswith("/"):
            layout_part = layout_part.lstrip("/")

        if layout_name not in layout_names:
            # Try part-based match using raw OpenXML inventory.
            if isinstance(layout_part, str) and layout_part in raw_by_part:
                # We'll validate placeholder_idx against raw placeholder idx set.
                raw_layout = raw_by_part[layout_part]
                raw_idxs = {
                    int(ph.get("idx"))
                    for ph in (raw_layout.get("placeholders") or [])
                    if isinstance(ph, dict) and isinstance(ph.get("idx"), int)
                }

                for slot in a.get("slots") or []:
                    if not isinstance(slot, dict):
                        errors.append(f"layout id '{aid}': slot entries must be objects")
                        continue
                    if "placeholder_idx" not in slot:
                        continue
                    idx = slot.get("placeholder_idx")
                    if not isinstance(idx, int):
                        errors.append(f"layout id '{aid}': slot '{slot.get('name','?')}' placeholder_idx must be int")
                        continue
                    if idx not in raw_idxs:
                        errors.append(
                            f"layout id '{aid}': layout_part '{layout_part}' missing placeholder idx={idx} for slot '{slot.get('name','?')}'"
                        )

                # If it's box-only, treat missing layout name as warning but still ok.
                if box_only:
                    errors.append(
                        f"warning: layout id '{aid}': slide layout name not found: '{layout_name}' (matched by layout_part)"
                    )
                continue

            if box_only:
                errors.append(
                    f"warning: layout id '{aid}': slide layout not found: '{layout_name}' (box-only definition; skipping layout checks)"
                )
                continue
            errors.append(f"layout id '{aid}': slide layout not found: '{layout_name}'")
            continue

        layout = layout_by_name(layout_name)
        if layout is None:
            if box_only:
                errors.append(
                    f"warning: layout id '{aid}': slide layout not found: '{layout_name}' (box-only definition; skipping layout checks)"
                )
                continue
            errors.append(f"layout id '{aid}': slide layout not found: '{layout_name}'")
            continue

        for slot in a.get("slots") or []:
            if not isinstance(slot, dict):
                errors.append(f"layout id '{aid}': slot entries must be objects")
                continue
            if "placeholder_idx" not in slot:
                continue
            idx = slot.get("placeholder_idx")
            if not isinstance(idx, int):
                errors.append(f"layout id '{aid}': slot '{slot.get('name','?')}' placeholder_idx must be int")
                continue
            found = False
            for ph in layout.placeholders:
                try:
                    if int(ph.placeholder_format.idx) == idx:
                        found = True
                        break
                except Exception:
                    continue
            if not found:
                errors.append(
                    f"layout id '{aid}': layout '{layout_name}' missing placeholder idx={idx} for slot '{slot.get('name','?')}'"
                )

    return errors


def validate_template(
    template_id: str,
    templates_dir: str | None = None,
    profile: str = "structural",
) -> TemplateValidationResult:
    errors: list[str] = []

    spec = load_template_spec(template_id, templates_dir=templates_dir)
    tdir = template_dir(template_id, templates_dir=templates_dir)
    pptx_path = tdir / "template.pptx"

    if profile not in {"structural", "standard", "extended"}:
        raise TemplateValidationError(f"Unknown validation profile: {profile}")

    archetypes = spec.get("archetypes") or []
    if not isinstance(archetypes, list) or not archetypes:
        return TemplateValidationResult(False, ["template spec must include non-empty 'archetypes' list"])

    # Template-native layout definitions (optional). These are additional layout definitions
    # that may be referenced directly by callers (namespaced IDs recommended).
    native = spec.get("native") or {}
    native_archetypes = []
    if isinstance(native, dict):
        native_archetypes = native.get("archetypes") or []

    if native_archetypes and not isinstance(native_archetypes, list):
        return TemplateValidationResult(False, ["template spec 'native.archetypes' must be a list when present"])

    deck = spec.get("deck") or {}

    # Semantic checks can run without a pptx.
    if profile in {"standard", "extended"}:
        errors.extend(_validate_semantic(archetypes, profile))
        if errors:
            return TemplateValidationResult(False, errors)


    has_pptx = pptx_path.exists()
    if not has_pptx and profile == "structural":
        return TemplateValidationResult(
            ok=True,
            errors=[f"warning: template.pptx not found; skipping layout/placeholder checks: {pptx_path}"],
        )

    if has_pptx:
        prs = Presentation(str(pptx_path))
        errors.extend(_validate_structural(archetypes, prs, pptx_path=str(pptx_path)))
        if native_archetypes:
            errors.extend(_validate_structural(native_archetypes, prs, pptx_path=str(pptx_path)))

    fatal = [e for e in errors if not str(e).startswith("warning:")]
    return TemplateValidationResult(ok=(len(fatal) == 0), errors=errors)
